"""Tests for the embedding service — pgvector RAG pipeline.

Uses mocked HTTP calls so no real API key is needed.
Tests that need pgvector in the DB are marked and skip gracefully
if the vector extension is unavailable.
"""

from __future__ import annotations

import os
from unittest.mock import AsyncMock, patch, MagicMock

import pytest
import pytest_asyncio
from sqlalchemy import text
from sqlalchemy.ext.asyncio import AsyncSession, async_sessionmaker, create_async_engine

from packages.core.models.base import Base
import packages.core.models  # noqa: F401

pytestmark = [pytest.mark.pgvector, pytest.mark.integration]

TEST_DATABASE_URL = os.getenv(
    "TEST_DATABASE_URL",
    "postgresql+asyncpg://manor:manor_secret@localhost:5434/manor_test",
)

EMBEDDING_DIMENSIONS = int(os.getenv("EMBEDDING_DIMENSIONS", "1024") or 1024)
FAKE_EMBEDDING = [0.01 * i for i in range(EMBEDDING_DIMENSIONS)]


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------


@pytest_asyncio.fixture
async def db_session():
    """Async DB session with vector extension enabled (if available)."""
    engine = create_async_engine(TEST_DATABASE_URL, echo=False)
    session_factory = async_sessionmaker(engine, class_=AsyncSession, expire_on_commit=False)

    has_vector = False
    try:
        async with engine.begin() as conn:
            await conn.execute(text("CREATE EXTENSION IF NOT EXISTS vector"))
            has_vector = True
    except Exception:
        pass

    async with engine.begin() as conn:
        await conn.run_sync(Base.metadata.drop_all)
        if has_vector:
            await conn.execute(text("CREATE EXTENSION IF NOT EXISTS vector"))
        await conn.run_sync(Base.metadata.create_all)
        if has_vector:
            await conn.execute(
                text(f"ALTER TABLE documents ADD COLUMN IF NOT EXISTS embedding vector({EMBEDDING_DIMENSIONS})")
            )

    async with session_factory() as session:
        yield session, has_vector

    await engine.dispose()


def _mock_embedding_response(embeddings: list[list[float]] | None = None):
    """Build a mock httpx response for the /embeddings endpoint."""
    if embeddings is None:
        embeddings = [FAKE_EMBEDDING]
    data = [{"embedding": emb, "index": i} for i, emb in enumerate(embeddings)]
    return MagicMock(
        status_code=200,
        json=lambda: {"data": data, "model": "text-embedding-3-small"},
        raise_for_status=lambda: None,
    )


def _set_manor_fs_root(monkeypatch, path) -> None:
    from packages.core.config import get_settings

    monkeypatch.setenv("MANOR_FS_ROOT", str(path))
    monkeypatch.setattr(get_settings(), "MANOR_FS_ROOT", str(path))


# ---------------------------------------------------------------------------
# Tests
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_generate_embedding_mock():
    """Mock httpx to return a fake embedding, verify configured vector size."""
    from packages.core.services.embedding_service import generate_embedding

    mock_resp = _mock_embedding_response()

    with patch.dict(os.environ, {"EMBEDDING_API_KEY": "test-key"}):
        with patch("packages.core.services.embedding_service.httpx.AsyncClient") as mock_cls:
            mock_client = AsyncMock()
            mock_client.post.return_value = mock_resp
            mock_client.__aenter__ = AsyncMock(return_value=mock_client)
            mock_client.__aexit__ = AsyncMock(return_value=False)
            mock_cls.return_value = mock_client

            result = await generate_embedding("Hello world")

    assert isinstance(result, list)
    assert len(result) == EMBEDDING_DIMENSIONS
    assert result == FAKE_EMBEDDING


@pytest.mark.asyncio
async def test_index_document_mock(db_session):
    """Mock embedding API, create a document, index it, verify vector_status='ready'."""
    session, has_vector = db_session
    if not has_vector:
        pytest.skip("pgvector extension not available")

    from packages.core.models.base import generate_ulid
    from packages.core.models.document import Document
    from packages.core.services.embedding_service import index_document

    doc_id = generate_ulid()
    doc = Document(
        id=doc_id,
        entity_id="ent_test",
        name="Test document about quantum computing",
        vector_status="pending",
    )
    session.add(doc)
    await session.flush()

    mock_resp = _mock_embedding_response()

    with patch.dict(os.environ, {"EMBEDDING_API_KEY": "test-key"}):
        with patch("packages.core.services.embedding_service.httpx.AsyncClient") as mock_cls:
            mock_client = AsyncMock()
            mock_client.post.return_value = mock_resp
            mock_client.__aenter__ = AsyncMock(return_value=mock_client)
            mock_client.__aexit__ = AsyncMock(return_value=False)
            mock_cls.return_value = mock_client

            success = await index_document(session, doc_id)

    assert success is True
    await session.commit()

    # Verify status
    row = await session.execute(
        text("SELECT vector_status FROM documents WHERE id = :id"),
        {"id": doc_id},
    )
    assert row.scalar_one() == "ready"


@pytest.mark.asyncio
async def test_index_document_missing_fs_path_is_skipped_not_ready_or_failed(
    db_session,
    tmp_path,
    monkeypatch,
):
    """Filesystem-backed docs cannot be marked ready when the real file is gone."""
    session, _has_vector = db_session

    from packages.core.models.base import generate_ulid
    from packages.core.models.document import Document
    from packages.core.services import embedding_service
    from packages.core.services.embedding_service import index_document

    entity_id = generate_ulid()
    (tmp_path / entity_id).mkdir()
    _set_manor_fs_root(monkeypatch, tmp_path)

    doc_id = generate_ulid()
    doc = Document(
        id=doc_id,
        entity_id=entity_id,
        name="missing.md",
        fs_path="missing.md",
        file_type="md",
        vector_status="pending",
    )
    session.add(doc)
    await session.flush()

    resolver = AsyncMock(return_value=None)
    monkeypatch.setattr(embedding_service, "_resolve_embedding_config", resolver)
    with patch("packages.core.services.embedding_service.httpx.AsyncClient") as mock_cls:
        success = await index_document(session, doc_id)

    assert success is False
    resolver.assert_not_called()
    mock_cls.assert_not_called()
    await session.refresh(doc)
    assert doc.vector_status == "skipped"
    integrity = (doc.metadata_ or {}).get("file_integrity") or {}
    assert integrity["status"] == "missing"
    assert integrity["fs_path"] == "missing.md"


@pytest.mark.asyncio
async def test_read_document_content_missing_fs_path_does_not_use_metadata_fallback(tmp_path, monkeypatch):
    """A stale fs_path must not be embedded from metadata/name fallback text."""
    from packages.core.models.base import generate_ulid
    from packages.core.models.document import Document
    from packages.core.services.embedding_service import DocumentContentUnavailable, _read_document_content

    entity_id = generate_ulid()
    (tmp_path / entity_id).mkdir()
    _set_manor_fs_root(monkeypatch, tmp_path)

    doc = Document(
        id=generate_ulid(),
        entity_id=entity_id,
        name="missing.md",
        fs_path="missing.md",
        file_type="md",
        metadata_={"content": "metadata fallback must not be indexed"},
        vector_status="pending",
    )

    with pytest.raises(DocumentContentUnavailable) as exc_info:
        await _read_document_content(doc)

    assert exc_info.value.status == "missing"


@pytest.mark.asyncio
async def test_index_document_reads_xlsx_extracted_text_beyond_500_rows(db_session, tmp_path, monkeypatch):
    """RAG indexing should embed extracted spreadsheet text, not raw ZIP bytes or a tiny preview."""
    session, has_vector = db_session
    if not has_vector:
        pytest.skip("pgvector extension not available")

    from openpyxl import Workbook
    from packages.core.models.base import generate_ulid
    from packages.core.models.document import Document
    from packages.core.services.embedding_service import index_document

    entity_id = generate_ulid()
    entity_root = tmp_path / entity_id
    entity_root.mkdir()
    rel_path = "Manor_AI_功能开发验证列表.xlsx"
    workbook_path = entity_root / rel_path
    wb = Workbook()
    ws = wb.active
    ws.title = "验证列表"
    ws.append(["编号", "功能", "状态", "验证备注"])
    for index in range(1, 651):
        ws.append([index, f"功能 {index}", "待验证", f"第 {index} 行备注"])
    wb.save(workbook_path)

    _set_manor_fs_root(monkeypatch, tmp_path)

    doc_id = generate_ulid()
    session.add(
        Document(
            id=doc_id,
            entity_id=entity_id,
            name=rel_path,
            fs_path=rel_path,
            file_type="xlsx",
            mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            vector_status="pending",
        )
    )
    await session.flush()

    captured_inputs: list[str] = []
    mock_resp = _mock_embedding_response()

    with patch.dict(os.environ, {"EMBEDDING_API_KEY": "test-key"}):
        with patch("packages.core.services.embedding_service.httpx.AsyncClient") as mock_cls:
            mock_client = AsyncMock()

            async def _post(_url, json=None, headers=None):
                value = (json or {}).get("input")
                if isinstance(value, list):
                    captured_inputs.extend(value)
                elif isinstance(value, str):
                    captured_inputs.append(value)
                return mock_resp

            mock_client.post.side_effect = _post
            mock_client.__aenter__ = AsyncMock(return_value=mock_client)
            mock_client.__aexit__ = AsyncMock(return_value=False)
            mock_cls.return_value = mock_client

            success = await index_document(session, doc_id)

    assert success is True
    indexed_text = "\n".join(captured_inputs)
    assert "[Sheet: 验证列表]" in indexed_text
    assert "650 | 功能 650 | 待验证 | 第 650 行备注" in indexed_text
    assert "PK\x03\x04" not in indexed_text


@pytest.mark.asyncio
async def test_index_document_reads_docx_and_pptx_extracted_text(db_session, tmp_path, monkeypatch):
    """RAG indexing should extract Word/PPT text before embedding."""
    session, has_vector = db_session
    if not has_vector:
        pytest.skip("pgvector extension not available")

    docx = pytest.importorskip("docx")
    pptx = pytest.importorskip("pptx")
    DocumentBuilder = docx.Document
    Presentation = pptx.Presentation
    from pptx.util import Inches
    from packages.core.models.base import generate_ulid
    from packages.core.models.document import Document
    from packages.core.services.embedding_service import index_document

    entity_id = generate_ulid()
    entity_root = tmp_path / entity_id
    entity_root.mkdir()

    docx_rel = "guest-brief.docx"
    docx_path = entity_root / docx_rel
    word = DocumentBuilder()
    word.add_paragraph("Guest messaging status: Ready for guest replies.")
    word.save(docx_path)

    pptx_rel = "owner-update.pptx"
    pptx_path = entity_root / pptx_rel
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(8), Inches(1))
    box.text = "AI edit status: Client-ready narrative"
    deck.save(pptx_path)

    _set_manor_fs_root(monkeypatch, tmp_path)

    word_doc_id = generate_ulid()
    deck_doc_id = generate_ulid()
    session.add_all(
        [
            Document(
                id=word_doc_id,
                entity_id=entity_id,
                name=docx_rel,
                fs_path=docx_rel,
                file_type="docx",
                mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                vector_status="pending",
            ),
            Document(
                id=deck_doc_id,
                entity_id=entity_id,
                name=pptx_rel,
                fs_path=pptx_rel,
                file_type="pptx",
                mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                vector_status="pending",
            ),
        ]
    )
    await session.flush()

    captured_inputs: list[str] = []
    mock_resp = _mock_embedding_response()

    with patch.dict(os.environ, {"EMBEDDING_API_KEY": "test-key"}):
        with patch("packages.core.services.embedding_service.httpx.AsyncClient") as mock_cls:
            mock_client = AsyncMock()

            async def _post(_url, json=None, headers=None):
                value = (json or {}).get("input")
                if isinstance(value, list):
                    captured_inputs.extend(value)
                elif isinstance(value, str):
                    captured_inputs.append(value)
                return mock_resp

            mock_client.post.side_effect = _post
            mock_client.__aenter__ = AsyncMock(return_value=mock_client)
            mock_client.__aexit__ = AsyncMock(return_value=False)
            mock_cls.return_value = mock_client

            assert await index_document(session, word_doc_id) is True
            assert await index_document(session, deck_doc_id) is True

    indexed_text = "\n".join(captured_inputs)
    assert "Ready for guest replies" in indexed_text
    assert "Client-ready narrative" in indexed_text
    assert "PK\x03\x04" not in indexed_text


@pytest.mark.asyncio
async def test_search_similar_mock(db_session):
    """Insert a doc with a known embedding, search with similar vector."""
    session, has_vector = db_session
    if not has_vector:
        pytest.skip("pgvector extension not available")

    from packages.core.models.base import generate_ulid
    from packages.core.models.document import Document
    from packages.core.services.embedding_service import search_similar

    doc_id = generate_ulid()
    doc = Document(
        id=doc_id,
        entity_id="ent_search",
        name="Machine learning basics",
        vector_status="ready",
    )
    session.add(doc)
    await session.flush()

    # Insert embedding via raw SQL
    vec_str = "[" + ",".join(str(v) for v in FAKE_EMBEDDING) + "]"
    await session.execute(
        text("UPDATE documents SET embedding = CAST(:vec AS vector) WHERE id = :id"),
        {"vec": vec_str, "id": doc_id},
    )
    await session.commit()

    # Search with the same embedding — should get a perfect match
    mock_resp = _mock_embedding_response()

    with patch.dict(os.environ, {"EMBEDDING_API_KEY": "test-key"}):
        with patch("packages.core.services.embedding_service.httpx.AsyncClient") as mock_cls:
            mock_client = AsyncMock()
            mock_client.post.return_value = mock_resp
            mock_client.__aenter__ = AsyncMock(return_value=mock_client)
            mock_client.__aexit__ = AsyncMock(return_value=False)
            mock_cls.return_value = mock_client

            results = await search_similar(session, "ent_search", "what is machine learning", threshold=0.5)

    assert len(results) >= 1
    assert results[0]["document_id"] == doc_id
    assert results[0]["score"] >= 0.99  # same vector = ~1.0 similarity


@pytest.mark.asyncio
async def test_hybrid_search_mock(db_session):
    """Test combined vector + text results."""
    session, has_vector = db_session
    if not has_vector:
        pytest.skip("pgvector extension not available")

    from packages.core.models.base import generate_ulid
    from packages.core.models.document import Document
    from packages.core.services.embedding_service import hybrid_search

    # Create two docs — one with embedding, one without
    doc1_id = generate_ulid()
    doc1 = Document(
        id=doc1_id,
        entity_id="ent_hybrid",
        name="Python programming guide",
        vector_status="ready",
    )
    doc2_id = generate_ulid()
    doc2 = Document(
        id=doc2_id,
        entity_id="ent_hybrid",
        name="Python cookbook recipes",
        vector_status="pending",
    )
    session.add_all([doc1, doc2])
    await session.flush()

    vec_str = "[" + ",".join(str(v) for v in FAKE_EMBEDDING) + "]"
    await session.execute(
        text("UPDATE documents SET embedding = CAST(:vec AS vector) WHERE id = :id"),
        {"vec": vec_str, "id": doc1_id},
    )
    await session.commit()

    mock_resp = _mock_embedding_response()

    with patch.dict(os.environ, {"EMBEDDING_API_KEY": "test-key"}):
        with patch("packages.core.services.embedding_service.httpx.AsyncClient") as mock_cls:
            mock_client = AsyncMock()
            mock_client.post.return_value = mock_resp
            mock_client.__aenter__ = AsyncMock(return_value=mock_client)
            mock_client.__aexit__ = AsyncMock(return_value=False)
            mock_cls.return_value = mock_client

            results = await hybrid_search(session, "ent_hybrid", "Python")

    doc_ids = [r["document_id"] for r in results]
    # Both docs should appear — one from vector, one from text
    assert doc1_id in doc_ids
    assert doc2_id in doc_ids


@pytest.mark.asyncio
async def test_no_api_key_fallback(db_session, tmp_path, monkeypatch):
    """When no EMBEDDING_API_KEY, hybrid_search falls back to text-only."""
    session, has_vector = db_session

    from packages.core.models.base import generate_ulid
    from packages.core.models.document import Document
    from packages.core.services.embedding_service import hybrid_search

    doc_id = generate_ulid()
    entity_root = tmp_path / "ent_nokey"
    entity_root.mkdir()
    (entity_root / "fallback.md").write_text("Fallback document body from the filesystem.", encoding="utf-8")

    doc = Document(
        id=doc_id,
        entity_id="ent_nokey",
        name="Fallback document test",
        fs_path="fallback.md",
        vector_status="pending",
    )
    session.add(doc)
    await session.flush()
    await session.commit()

    _set_manor_fs_root(monkeypatch, tmp_path)

    # Ensure no embedding keys are set
    env_patch = {
        "EMBEDDING_API_KEY": "",
        "OPENAI_API_KEY": "",
        "MANOR_FS_ROOT": str(tmp_path),
    }
    with patch.dict(os.environ, env_patch, clear=False):
        results = await hybrid_search(session, "ent_nokey", "Fallback")

    assert len(results) >= 1
    assert results[0]["document_id"] == doc_id
    # Score should come from the text/lexical fallback path.
    assert results[0]["score"] >= 0.5
    assert "Fallback document body" in results[0]["content_preview"]


@pytest.mark.asyncio
async def test_vector_type_column(db_session):
    """Verify the VectorType column can store and retrieve embeddings via raw SQL."""
    session, has_vector = db_session
    if not has_vector:
        pytest.skip("pgvector extension not available")

    from packages.core.models.base import generate_ulid
    from packages.core.models.document import Document

    doc_id = generate_ulid()
    doc = Document(
        id=doc_id,
        entity_id="ent_vtype",
        name="Vector type test",
        vector_status="pending",
    )
    session.add(doc)
    await session.flush()

    # Store embedding via raw SQL
    vec_str = "[" + ",".join(str(float(i)) for i in range(EMBEDDING_DIMENSIONS)) + "]"
    await session.execute(
        text("UPDATE documents SET embedding = CAST(:vec AS vector) WHERE id = :id"),
        {"vec": vec_str, "id": doc_id},
    )
    await session.commit()

    # Retrieve and verify the embedding column is accessible
    row = await session.execute(
        text("SELECT embedding::text FROM documents WHERE id = :id"),
        {"id": doc_id},
    )
    stored = row.scalar_one()
    assert stored is not None
    # Parse the stored vector string
    values = [float(x) for x in stored.strip("[]").split(",")]
    assert len(values) == EMBEDDING_DIMENSIONS
    assert values[0] == 0.0
    assert values[1] == 1.0
