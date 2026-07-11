from __future__ import annotations

import pytest

from packages.core.models.base import generate_ulid
from packages.core.models.document import Document, DocumentGroup, DocumentGroupMember, VectorStatus
from packages.core.services.document_metadata import merge_document_metadata
from packages.core.services.document_service import list_documents


@pytest.mark.asyncio
async def test_list_documents_workspace_filter_includes_group_and_artifact_provenance(db_session):
    entity_id = generate_ulid()
    workspace_id = generate_ulid()
    group_id = generate_ulid()

    group_doc = Document(
        id=generate_ulid(),
        entity_id=entity_id,
        name="brand-guide.md",
        file_type="md",
        source="upload",
        vector_status=VectorStatus.READY,
    )
    artifact_doc = Document(
        id=generate_ulid(),
        entity_id=entity_id,
        name="weekly-report.json",
        file_type="json",
        source="agent",
        vector_status=VectorStatus.READY,
        metadata_=merge_document_metadata(
            origin={"workspace_id": workspace_id, "task_id": "task_1"},
            artifact={"role": "final"},
        ),
    )
    other_doc = Document(
        id=generate_ulid(),
        entity_id=entity_id,
        name="other-workspace.json",
        file_type="json",
        source="agent",
        vector_status=VectorStatus.READY,
        metadata_=merge_document_metadata(
            origin={"workspace_id": generate_ulid()},
            artifact={"role": "final"},
        ),
    )
    group = DocumentGroup(
        id=group_id,
        entity_id=entity_id,
        workspace_id=workspace_id,
        name="Workspace Knowledge",
        settings={"default_collection": True},
    )

    db_session.add_all([group_doc, artifact_doc, other_doc, group])
    db_session.add(DocumentGroupMember(document_id=group_doc.id, group_id=group_id))
    await db_session.commit()

    docs, total = await list_documents(db_session, entity_id, workspace_id=workspace_id)

    ids = {doc.id for doc in docs}
    assert total == 2
    assert group_doc.id in ids
    assert artifact_doc.id in ids
    assert other_doc.id not in ids


@pytest.mark.asyncio
async def test_hybrid_search_group_ids_scope_to_knowledge_net(db_session, monkeypatch):
    from packages.core.services.embedding_service import hybrid_search

    entity_id = generate_ulid()
    net_a_id = generate_ulid()
    net_b_id = generate_ulid()
    bucket_id = generate_ulid()
    doc_a = Document(
        id=generate_ulid(),
        entity_id=entity_id,
        name="Lease Policy Alpha",
        file_type="md",
        source="upload",
        vector_status=VectorStatus.READY,
    )
    doc_b = Document(
        id=generate_ulid(),
        entity_id=entity_id,
        name="Lease Policy Beta",
        file_type="md",
        source="upload",
        vector_status=VectorStatus.READY,
    )
    db_session.add_all(
        [
            doc_a,
            doc_b,
            DocumentGroup(
                id=net_a_id,
                entity_id=entity_id,
                workspace_id=generate_ulid(),
                name="Lease Net A",
                settings={"kind": "knowledge_net"},
            ),
            DocumentGroup(
                id=net_b_id,
                entity_id=entity_id,
                workspace_id=generate_ulid(),
                name="Lease Net B",
                settings={"kind": "knowledge_net"},
            ),
            DocumentGroup(
                id=bucket_id,
                entity_id=entity_id,
                workspace_id=generate_ulid(),
                name="Workspace Files",
                settings={"workspace_file_bucket": True},
            ),
            DocumentGroupMember(document_id=doc_a.id, group_id=net_a_id),
            DocumentGroupMember(document_id=doc_b.id, group_id=net_b_id),
            DocumentGroupMember(document_id=doc_b.id, group_id=bucket_id),
        ]
    )
    await db_session.commit()

    monkeypatch.setenv("EMBEDDING_API_KEY", "")
    monkeypatch.setenv("OPENAI_API_KEY", "")

    scoped = await hybrid_search(db_session, entity_id, "Lease Policy", group_ids=[net_a_id])
    bucket_scoped = await hybrid_search(db_session, entity_id, "Lease Policy", group_ids=[bucket_id])

    assert [result["document_id"] for result in scoped] == [doc_a.id]
    assert bucket_scoped == []


@pytest.mark.asyncio
async def test_hybrid_search_scoped_knowledge_net_reads_content_without_embeddings(db_session, monkeypatch):
    from packages.core.services.embedding_service import hybrid_search

    entity_id = generate_ulid()
    net_id = generate_ulid()
    doc = Document(
        id=generate_ulid(),
        entity_id=entity_id,
        name="Brand Voice XHS X Safety.md",
        file_type="md",
        source="upload",
        vector_status=VectorStatus.READY,
        metadata_={
            "content_text": (
                "Tone: practical, founder-operated, low-hype. Public posts on "
                "X/Xiaohongshu must stay as drafts until the user approves."
            )
        },
    )
    db_session.add_all(
        [
            doc,
            DocumentGroup(
                id=net_id,
                entity_id=entity_id,
                workspace_id=generate_ulid(),
                name="Brand Voice & Safety",
                settings={"kind": "knowledge_net"},
            ),
            DocumentGroupMember(document_id=doc.id, group_id=net_id),
        ]
    )
    await db_session.commit()

    monkeypatch.setenv("EMBEDDING_API_KEY", "")
    monkeypatch.setenv("OPENAI_API_KEY", "")

    results = await hybrid_search(
        db_session,
        entity_id,
        "approval rules for public posts on X and Xiaohongshu",
        group_ids=[net_id],
    )

    assert [result["document_id"] for result in results] == [doc.id]
    assert "must stay as drafts" in results[0]["content_preview"]


@pytest.mark.asyncio
async def test_spreadsheet_content_preview_is_long_enough_for_status_rows(tmp_path, monkeypatch):
    from openpyxl import Workbook
    from packages.core.services.embedding_service import _build_content_preview

    entity_id = generate_ulid()
    entity_root = tmp_path / entity_id
    entity_root.mkdir()
    rel_path = "Manor_AI_功能开发验证列表.xlsx"
    workbook_path = entity_root / rel_path
    wb = Workbook()
    ws = wb.active
    ws.title = "验证列表"
    ws.append(["编号", "功能", "状态", "验证备注"])
    for index in range(1, 90):
        ws.append([index, f"功能 {index}", "待验证", f"第 {index} 行备注"])
    wb.save(workbook_path)

    monkeypatch.setenv("MANOR_FS_ROOT", str(tmp_path))

    preview = await _build_content_preview(
        entity_id=entity_id,
        fs_path=rel_path,
        mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        file_type="xlsx",
        metadata={},
        name=rel_path,
        max_chars=500,
    )

    assert len(preview) > 500
    assert "80 | 功能 80 | 待验证 | 第 80 行备注" in preview
    assert "PK\x03\x04" not in preview


@pytest.mark.asyncio
async def test_scoped_lexical_search_reads_deep_spreadsheet_rows_without_embeddings(db_session, tmp_path, monkeypatch):
    from openpyxl import Workbook
    from packages.core.services.embedding_service import hybrid_search

    entity_id = generate_ulid()
    net_id = generate_ulid()
    entity_root = tmp_path / entity_id
    entity_root.mkdir()
    rel_path = "Manor_AI_功能开发验证列表.xlsx"
    workbook_path = entity_root / rel_path
    wb = Workbook()
    ws = wb.active
    ws.title = "验证列表"
    ws.append(["编号", "功能", "状态", "验证备注"])
    for index in range(1, 651):
        ws.append([index, f"功能 {index}", "待验证", f"第 {index} 行备注"])
    wb.save(workbook_path)

    doc = Document(
        id=generate_ulid(),
        entity_id=entity_id,
        name=rel_path,
        fs_path=rel_path,
        file_type="xlsx",
        mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        source="upload",
        vector_status=VectorStatus.READY,
    )
    db_session.add_all(
        [
            doc,
            DocumentGroup(
                id=net_id,
                entity_id=entity_id,
                workspace_id=generate_ulid(),
                name="开发验证列表",
                settings={"kind": "knowledge_net"},
            ),
            DocumentGroupMember(document_id=doc.id, group_id=net_id),
        ]
    )
    await db_session.commit()

    monkeypatch.setenv("MANOR_FS_ROOT", str(tmp_path))
    monkeypatch.setenv("EMBEDDING_API_KEY", "")
    monkeypatch.setenv("OPENAI_API_KEY", "")

    results = await hybrid_search(
        db_session,
        entity_id,
        "功能 650 第 650 行备注",
        group_ids=[net_id],
    )

    assert [result["document_id"] for result in results] == [doc.id]
    assert "650 | 功能 650 | 待验证 | 第 650 行备注" in results[0]["content_preview"]
    assert "PK\x03\x04" not in results[0]["content_preview"]


@pytest.mark.asyncio
async def test_office_content_previews_are_extracted_text(tmp_path, monkeypatch):
    docx = pytest.importorskip("docx")
    pptx = pytest.importorskip("pptx")
    DocumentBuilder = docx.Document
    Presentation = pptx.Presentation
    from pptx.util import Inches
    from packages.core.services.embedding_service import _build_content_preview

    entity_id = generate_ulid()
    entity_root = tmp_path / entity_id
    entity_root.mkdir()

    docx_rel = "guest-brief.docx"
    docx_path = entity_root / docx_rel
    word = DocumentBuilder()
    word.add_paragraph("Guest messaging status: Ready for guest replies.")
    word.save(docx_path)

    pptx_rel = "owner-update.pptx"
    pptx_path = entity_root / pptx_rel
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(8), Inches(1))
    box.text = "AI edit status: Client-ready narrative"
    deck.save(pptx_path)

    monkeypatch.setenv("MANOR_FS_ROOT", str(tmp_path))

    docx_preview = await _build_content_preview(
        entity_id=entity_id,
        fs_path=docx_rel,
        mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        file_type="docx",
        metadata={},
        name=docx_rel,
    )
    pptx_preview = await _build_content_preview(
        entity_id=entity_id,
        fs_path=pptx_rel,
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        file_type="pptx",
        metadata={},
        name=pptx_rel,
    )

    assert "Ready for guest replies" in docx_preview
    assert "Client-ready narrative" in pptx_preview
    assert "PK\x03\x04" not in docx_preview
    assert "PK\x03\x04" not in pptx_preview
