import sys
from pathlib import Path

import pytest

ROOT = Path(__file__).resolve().parents[1]
SCRIPTS_DIR = ROOT / "packages/core/ai/skills/pptx/scripts"
sys.path.insert(0, str(SCRIPTS_DIR))

import images_to_pptx  # noqa: E402

PIL = pytest.importorskip("PIL")
from PIL import Image  # noqa: E402

pptx = pytest.importorskip("pptx")
from pptx import Presentation  # noqa: E402


def _make_png(path: Path, width: int, height: int, color=(200, 30, 30)) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    Image.new("RGB", (width, height), color).save(path)


def _project(tmp_path: Path, name: str = "demo_ppt169_20260615") -> Path:
    project = tmp_path / name
    (project / "images").mkdir(parents=True)
    (project / "exports").mkdir(parents=True)
    return project


def test_collect_images_natural_order(tmp_path: Path) -> None:
    images_dir = tmp_path / "images"
    for i in (1, 2, 10):
        _make_png(images_dir / f"page_{i}.png", 16, 9)
    # A non-page image must be ignored by the default glob.
    _make_png(images_dir / "logo.png", 16, 9)

    found = images_to_pptx.collect_images(images_dir, "page_*")

    assert [p.name for p in found] == ["page_1.png", "page_2.png", "page_10.png"]


def test_build_pptx_one_slide_per_image_and_slide_size(tmp_path: Path) -> None:
    project = _project(tmp_path)
    images_dir = project / "images"
    _make_png(images_dir / "page_01.png", 1280, 720)
    _make_png(images_dir / "page_02.png", 1280, 720)

    images = images_to_pptx.collect_images(images_dir, "page_*")
    output = project / "exports" / "demo.pptx"
    images_to_pptx.build_pptx(images, 1280, 720, output, fit="cover")

    assert output.is_file()
    prs = Presentation(str(output))
    assert len(prs.slides) == 2
    # 1280 px at 96 dpi -> 1280/96 inch -> EMU
    assert prs.slide_width == int(round(1280 * images_to_pptx.EMU_PER_PIXEL))
    assert prs.slide_height == int(round(720 * images_to_pptx.EMU_PER_PIXEL))
    # Each slide carries exactly one picture.
    for slide in prs.slides:
        pics = [s for s in slide.shapes if s.shape_type == 13]  # PICTURE
        assert len(pics) == 1


def test_cover_crops_mismatched_aspect(tmp_path: Path) -> None:
    project = _project(tmp_path)
    images_dir = project / "images"
    # A square image into a 16:9 slide must be cropped top/bottom under cover.
    _make_png(images_dir / "page_01.png", 1000, 1000)

    images = images_to_pptx.collect_images(images_dir, "page_*")
    output = project / "exports" / "demo.pptx"
    images_to_pptx.build_pptx(images, 1280, 720, output, fit="cover")

    prs = Presentation(str(output))
    pic = next(s for s in prs.slides[0].shapes if s.shape_type == 13)
    assert pic.crop_top > 0
    assert pic.crop_bottom > 0
    assert pic.crop_left == 0
    assert pic.crop_right == 0


def test_default_fit_contains_mismatched_aspect(tmp_path: Path, monkeypatch) -> None:
    project = _project(tmp_path)
    images_dir = project / "images"
    # A square page image should be fully visible by default, not center-cropped.
    _make_png(images_dir / "page_01.png", 1000, 1000)
    output = project / "exports" / "deck.pptx"

    monkeypatch.setattr(
        "sys.argv",
        ["images_to_pptx.py", str(project), "-o", str(output)],
    )
    rc = images_to_pptx.main()

    assert rc == 0
    prs = Presentation(str(output))
    pic = next(s for s in prs.slides[0].shapes if s.shape_type == 13)
    assert pic.crop_top == 0
    assert pic.crop_bottom == 0
    assert pic.crop_left == 0
    assert pic.crop_right == 0

    slide_h = int(round(720 * images_to_pptx.EMU_PER_PIXEL))
    slide_w = int(round(1280 * images_to_pptx.EMU_PER_PIXEL))
    assert pic.width == slide_h
    assert pic.height == slide_h
    assert pic.left == (slide_w - slide_h) // 2
    assert pic.top == 0


def test_resolve_format_from_project_name(tmp_path: Path) -> None:
    project = _project(tmp_path, name="myslides_story_20260615")
    assert images_to_pptx.resolve_format(project, None) == "story"
    # Explicit flag wins.
    assert images_to_pptx.resolve_format(project, "ppt43") == "ppt43"
    # Unknown directory name falls back to ppt169.
    other = _project(tmp_path, name="loose-folder")
    assert images_to_pptx.resolve_format(other, None) == "ppt169"


def test_slide_pixels_known_and_default() -> None:
    assert images_to_pptx.slide_pixels("ppt169") == (1280, 720)
    assert images_to_pptx.slide_pixels("story") == (1080, 1920)
    assert images_to_pptx.slide_pixels("does-not-exist") == (1280, 720)


def test_cleanup_removes_source_images_keeps_deck(tmp_path: Path, monkeypatch) -> None:
    project = _project(tmp_path)
    images_dir = project / "images"
    _make_png(images_dir / "page_01.png", 1280, 720)
    _make_png(images_dir / "page_02.png", 1280, 720)
    output = project / "exports" / "deck.pptx"

    monkeypatch.setattr(
        "sys.argv",
        ["images_to_pptx.py", str(project), "--cleanup", "-o", str(output)],
    )
    rc = images_to_pptx.main()

    assert rc == 0
    assert output.is_file()  # deck kept
    # Intermediate page images removed.
    assert not (images_dir / "page_01.png").exists()
    assert not (images_dir / "page_02.png").exists()


def test_no_cleanup_keeps_source_images(tmp_path: Path, monkeypatch) -> None:
    project = _project(tmp_path)
    images_dir = project / "images"
    _make_png(images_dir / "page_01.png", 1280, 720)
    output = project / "exports" / "deck.pptx"

    monkeypatch.setattr(
        "sys.argv",
        ["images_to_pptx.py", str(project), "-o", str(output)],
    )
    rc = images_to_pptx.main()

    assert rc == 0
    assert (images_dir / "page_01.png").exists()  # kept without --cleanup
