from __future__ import annotations

import json
from types import SimpleNamespace
from unittest.mock import AsyncMock

import pytest

from packages.core.ai.tools import file_tools as file_tools_module
from packages.core.ai.tools.document_tools import _cache_key as document_cache_key
from packages.core.ai.tools.document_tools import _doc_to_dict
from packages.core.ai.tools.document_tools import _generate_document_file
from packages.core.ai.tools.file_tools import (
    _delete_file,
    _edit_file,
    _glob_files,
    _grep_files,
    _list_files,
    _read_file,
    _write_file,
)
from packages.core.ai.tools.extended_tools import _web_fetch_handler
from packages.core.ai.tools.generate_file_tool import GENERATE_FILE_SCHEMA
from packages.core.ai.runtime.manor_actions import (
    _runtime_manor_invalidate_read_cache as _invalidate_read_cache,
    _runtime_manor_read_cache_key as _read_cache_key,
    _runtime_staff_summary as _staff_summary,
)
from packages.core.ai.tools.manor_tool import MANOR_SCHEMA, _doc_summary
from packages.core.ai.tools.sandbox_tools import _SANDBOX_SAVE_RESULT_SCHEMA, _sandbox_read_file
from packages.core.ai.tools.bash_tool import _execute_local
from packages.core.ai.tool_pool import ToolPool
from packages.core.ai.runtime.tool_registry import runtime_registered_tool_surface_from_schemas
from packages.core.config import get_settings


def _minimal_schema(name: str) -> dict:
    return {
        "type": "function",
        "function": {
            "name": name,
            "description": name,
            "parameters": {"type": "object", "properties": {}},
        },
    }


def test_file_tools_expose_single_edit_tool_for_text_and_spreadsheets():
    names = [schema["function"]["name"] for schema, _handler in file_tools_module.get_tools()]
    edit_schema = next(
        schema for schema, _handler in file_tools_module.get_tools() if schema["function"]["name"] == "edit_file"
    )
    props = edit_schema["function"]["parameters"]["properties"]

    assert names.count("edit_file") == 1
    assert "edit_spreadsheet" not in names
    assert props["operation"]["enum"] == ["replace_text", "set_cell", "update_row", "append_row"]
    assert {"sheet", "cell", "match_column", "updates", "row"} <= set(props)


def _tool_surface_for_pool(pool: ToolPool, *, is_master: bool = False) -> tuple[list[dict], list[str]]:
    surface = runtime_registered_tool_surface_from_schemas(
        pool.registered_tool_schemas(),
        is_master=is_master,
    )
    return list(surface.prompt_schemas), list(surface.deferred_tool_names)


@pytest.mark.asyncio
async def test_read_file_returns_bounded_slice_with_continuation(tmp_path):
    settings = get_settings()
    old_enabled = settings.MANOR_FS_ENABLED
    old_root = settings.MANOR_FS_ROOT
    settings.MANOR_FS_ENABLED = True
    settings.MANOR_FS_ROOT = str(tmp_path)
    try:
        entity_id = "ent_1"
        entity_root = tmp_path / entity_id
        entity_root.mkdir()
        content = "".join(f"{i:03d} " + ("x" * 180) + "\n" for i in range(20))
        (entity_root / "big.txt").write_text(content, encoding="utf-8")

        result = json.loads(
            await _read_file(
                entity_id=entity_id,
                path="big.txt",
                limit=10,
                max_chars=1000,
            )
        )

        assert len(result["content"]) <= 1000
        assert result["lines_returned"] == 5
        assert result["next_offset"] == 5
        assert result["next_char_offset"] == len(result["content"])
        assert result["truncated"] is True
        assert len(result["source_sha256"]) == 64
        assert len(result["slice_sha256"]) == 64
        assert result["size"] == len(content.encode("utf-8"))
        assert isinstance(result["mtime_ns"], int)
        assert "char_offset=next_char_offset" in result["hint"]

        continuation = json.loads(
            await _read_file(
                entity_id=entity_id,
                path="big.txt",
                char_offset=result["next_char_offset"],
                limit=10,
                expected_sha256=result["source_sha256"],
            )
        )
        assert continuation["mode"] == "char"
        assert continuation["char_offset"] == result["next_char_offset"]
        assert continuation["source_sha256"] == result["source_sha256"]
    finally:
        settings.MANOR_FS_ENABLED = old_enabled
        settings.MANOR_FS_ROOT = old_root


@pytest.mark.asyncio
async def test_read_file_exact_slices_reconstruct_source(tmp_path):
    settings = get_settings()
    old_enabled = settings.MANOR_FS_ENABLED
    old_root = settings.MANOR_FS_ROOT
    settings.MANOR_FS_ENABLED = True
    settings.MANOR_FS_ROOT = str(tmp_path)
    try:
        entity_id = "ent_1"
        entity_root = tmp_path / entity_id
        entity_root.mkdir()
        content = "alpha\n" + ("b" * 2500) + "\nomega\n"
        (entity_root / "long-line.txt").write_text(content, encoding="utf-8")

        chunks: list[str] = []
        next_char_offset: int | None = 0
        source_sha256 = None
        while next_char_offset is not None:
            result = json.loads(
                await _read_file(
                    entity_id=entity_id,
                    path="long-line.txt",
                    char_offset=next_char_offset,
                    max_chars=1000,
                    expected_sha256=source_sha256,
                )
            )
            source_sha256 = result["source_sha256"]
            chunks.append(result["content"])
            next_char_offset = result["next_char_offset"]

        assert "".join(chunks) == content
    finally:
        settings.MANOR_FS_ENABLED = old_enabled
        settings.MANOR_FS_ROOT = old_root


@pytest.mark.asyncio
async def test_read_file_extracts_xlsx_as_text(tmp_path):
    from openpyxl import Workbook

    settings = get_settings()
    old_enabled = settings.MANOR_FS_ENABLED
    old_root = settings.MANOR_FS_ROOT
    settings.MANOR_FS_ENABLED = True
    settings.MANOR_FS_ROOT = str(tmp_path)
    try:
        entity_id = "ent_1"
        entity_root = tmp_path / entity_id
        entity_root.mkdir()
        workbook_path = entity_root / "Manor_AI_功能开发验证列表.xlsx"
        wb = Workbook()
        ws = wb.active
        ws.title = "验证列表"
        ws.append(["编号", "功能", "状态", "验证备注"])
        for index in range(1, 81):
            ws.append([index, f"功能 {index}", "待验证", f"第 {index} 行备注"])
        wb.save(workbook_path)

        result = json.loads(
            await _read_file(
                entity_id=entity_id,
                path=workbook_path.name,
                limit=120,
                max_chars=20_000,
            )
        )

        assert result["total_lines"] == 82
        assert "[Sheet: 验证列表]" in result["content"]
        assert "80 | 功能 80 | 待验证 | 第 80 行备注" in result["content"]
        assert "PK\x03\x04" not in result["content"]
    finally:
        settings.MANOR_FS_ENABLED = old_enabled
        settings.MANOR_FS_ROOT = old_root


@pytest.mark.asyncio
async def test_edit_file_rejects_xlsx_without_corrupting_workbook(tmp_path):
    from openpyxl import Workbook, load_workbook

    settings = get_settings()
    old_enabled = settings.MANOR_FS_ENABLED
    old_root = settings.MANOR_FS_ROOT
    settings.MANOR_FS_ENABLED = True
    settings.MANOR_FS_ROOT = str(tmp_path)
    try:
        entity_id = "ent_1"
        entity_root = tmp_path / entity_id
        entity_root.mkdir()
        workbook_path = entity_root / "status.xlsx"
        wb = Workbook()
        ws = wb.active
        ws.append(["功能", "状态"])
        ws.append(["Public chat stream", "待验证"])
        wb.save(workbook_path)

        result = json.loads(
            await _edit_file(
                entity_id=entity_id,
                path=workbook_path.name,
                old_text="待验证",
                new_text="已验证",
            )
        )

        assert result["error"] == "unsupported_binary_edit"
        reopened = load_workbook(workbook_path)
        assert reopened.active["B2"].value == "待验证"
        reopened.close()
    finally:
        settings.MANOR_FS_ENABLED = old_enabled
        settings.MANOR_FS_ROOT = old_root


@pytest.mark.asyncio
async def test_edit_file_updates_matched_spreadsheet_status_row(tmp_path, monkeypatch):
    from openpyxl import Workbook, load_workbook

    settings = get_settings()
    old_enabled = settings.MANOR_FS_ENABLED
    old_root = settings.MANOR_FS_ROOT
    settings.MANOR_FS_ENABLED = True
    settings.MANOR_FS_ROOT = str(tmp_path)
    monkeypatch.setattr(
        "packages.core.services.ai_file_permissions.guard_ai_file_mutation",
        AsyncMock(return_value=None),
    )

    async def fake_sync_file_to_knowledge(**_kwargs):
        return SimpleNamespace(synced=True, document_id="doc_123", reason=None)

    monkeypatch.setattr(
        "packages.core.services.knowledge_sync.sync_file_to_knowledge",
        fake_sync_file_to_knowledge,
    )
    try:
        entity_id = "ent_1"
        entity_root = tmp_path / entity_id
        entity_root.mkdir()
        workbook_path = entity_root / "Manor_AI_功能开发验证列表.xlsx"
        wb = Workbook()
        ws = wb.active
        ws.title = "验证列表"
        ws.append(["编号", "功能", "状态", "验证备注"])
        ws.append([1, "Public chat stream", "待验证", ""])
        ws.append([2, "Channel language", "待验证", ""])
        wb.save(workbook_path)

        read_result = json.loads(await _read_file(entity_id=entity_id, path=workbook_path.name))
        result = json.loads(
            await _edit_file(
                entity_id=entity_id,
                path=workbook_path.name,
                sheet="验证列表",
                operation="update_row",
                header_row=1,
                match_column="功能",
                match_value="Public chat stream",
                updates={"状态": "已验证", "验证备注": "stream out and replying verified"},
                expected_sha256=read_result["source_sha256"],
            )
        )

        assert result["updated"] is True
        assert result["row_number"] == 2
        assert result["updated_cells"]["C2"]["old"] == "待验证"
        assert result["updated_cells"]["C2"]["new"] == "已验证"
        assert result["source_sha256"] != read_result["source_sha256"]
        assert result["knowledge_synced"] is True

        reopened = load_workbook(workbook_path)
        ws = reopened["验证列表"]
        assert ws["C2"].value == "已验证"
        assert ws["D2"].value == "stream out and replying verified"
        assert ws["C3"].value == "待验证"
        reopened.close()
    finally:
        settings.MANOR_FS_ENABLED = old_enabled
        settings.MANOR_FS_ROOT = old_root


@pytest.mark.asyncio
async def test_edit_file_appends_header_mapped_spreadsheet_row(tmp_path, monkeypatch):
    from openpyxl import Workbook, load_workbook

    settings = get_settings()
    old_enabled = settings.MANOR_FS_ENABLED
    old_root = settings.MANOR_FS_ROOT
    settings.MANOR_FS_ENABLED = True
    settings.MANOR_FS_ROOT = str(tmp_path)
    monkeypatch.setattr(
        "packages.core.services.ai_file_permissions.guard_ai_file_mutation",
        AsyncMock(return_value=None),
    )

    async def fake_sync_file_to_knowledge(**_kwargs):
        return SimpleNamespace(synced=True, document_id="doc_123", reason=None)

    monkeypatch.setattr(
        "packages.core.services.knowledge_sync.sync_file_to_knowledge",
        fake_sync_file_to_knowledge,
    )
    try:
        entity_id = "ent_1"
        entity_root = tmp_path / entity_id
        entity_root.mkdir()
        workbook_path = entity_root / "status.xlsx"
        wb = Workbook()
        ws = wb.active
        ws.append(["编号", "功能", "状态"])
        ws.append([1, "Existing", "已验证"])
        wb.save(workbook_path)

        result = json.loads(
            await _edit_file(
                entity_id=entity_id,
                path=workbook_path.name,
                operation="append_row",
                row={"编号": 2, "功能": "Excel update tool", "状态": "待验证"},
            )
        )

        assert result["updated"] is True
        assert result["row_number"] == 3
        reopened = load_workbook(workbook_path)
        ws = reopened.active
        assert [ws["A3"].value, ws["B3"].value, ws["C3"].value] == [2, "Excel update tool", "待验证"]
        reopened.close()
    finally:
        settings.MANOR_FS_ENABLED = old_enabled
        settings.MANOR_FS_ROOT = old_root


@pytest.mark.asyncio
async def test_real_user_edit_file_updates_docx_document(tmp_path, monkeypatch):
    docx = pytest.importorskip("docx")
    Document = docx.Document

    settings = get_settings()
    old_enabled = settings.MANOR_FS_ENABLED
    old_root = settings.MANOR_FS_ROOT
    settings.MANOR_FS_ENABLED = True
    settings.MANOR_FS_ROOT = str(tmp_path)
    monkeypatch.setattr(
        "packages.core.services.ai_file_permissions.guard_ai_file_mutation",
        AsyncMock(return_value=None),
    )

    async def fake_sync_file_to_knowledge(**_kwargs):
        return SimpleNamespace(synced=True, document_id="doc_word", reason=None)

    monkeypatch.setattr(
        "packages.core.services.knowledge_sync.sync_file_to_knowledge",
        fake_sync_file_to_knowledge,
    )
    try:
        entity_id = "ent_1"
        entity_root = tmp_path / entity_id
        entity_root.mkdir()
        document_path = entity_root / "guest-brief.docx"
        doc = Document()
        doc.add_paragraph("Guest messaging status: Draft for owner review.")
        table = doc.add_table(rows=1, cols=2)
        table.cell(0, 0).text = "Next step"
        table.cell(0, 1).text = "Send Draft to Das Blumenberg"
        doc.save(document_path)

        read_result = json.loads(await _read_file(entity_id=entity_id, path=document_path.name))
        assert "Draft for owner review" in read_result["content"]

        result = json.loads(
            await _edit_file(
                entity_id=entity_id,
                path=document_path.name,
                old_text="Draft for owner review",
                new_text="Ready for guest replies",
                expected_sha256=read_result["source_sha256"],
            )
        )

        assert result["edited"] is True
        assert result["file_type"] == "docx"
        assert result["replacements"] == 1
        assert result["knowledge_synced"] is True
        assert result["source_sha256"] != read_result["source_sha256"]

        reopened = Document(document_path)
        assert "Ready for guest replies" in "\n".join(p.text for p in reopened.paragraphs)
        assert (
            "Draft for owner review"
            not in json.loads(await _read_file(entity_id=entity_id, path=document_path.name))["content"]
        )
    finally:
        settings.MANOR_FS_ENABLED = old_enabled
        settings.MANOR_FS_ROOT = old_root


@pytest.mark.asyncio
async def test_real_user_edit_file_updates_pptx_deck(tmp_path, monkeypatch):
    pptx = pytest.importorskip("pptx")
    Presentation = pptx.Presentation
    from pptx.util import Inches

    settings = get_settings()
    old_enabled = settings.MANOR_FS_ENABLED
    old_root = settings.MANOR_FS_ROOT
    settings.MANOR_FS_ENABLED = True
    settings.MANOR_FS_ROOT = str(tmp_path)
    monkeypatch.setattr(
        "packages.core.services.ai_file_permissions.guard_ai_file_mutation",
        AsyncMock(return_value=None),
    )

    async def fake_sync_file_to_knowledge(**_kwargs):
        return SimpleNamespace(synced=True, document_id="doc_deck", reason=None)

    monkeypatch.setattr(
        "packages.core.services.knowledge_sync.sync_file_to_knowledge",
        fake_sync_file_to_knowledge,
    )
    try:
        entity_id = "ent_1"
        entity_root = tmp_path / entity_id
        entity_root.mkdir()
        deck_path = entity_root / "owner-update.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(8), Inches(0.7))
        title_box.text = "Das Blumenberg Owner Update"
        body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8), Inches(1.2))
        body_box.text = "AI edit status: Draft narrative"
        prs.save(deck_path)

        read_result = json.loads(await _read_file(entity_id=entity_id, path=deck_path.name))
        assert "Draft narrative" in read_result["content"]

        result = json.loads(
            await _edit_file(
                entity_id=entity_id,
                path=deck_path.name,
                old_text="Draft narrative",
                new_text="Client-ready narrative",
                expected_sha256=read_result["source_sha256"],
            )
        )

        assert result["edited"] is True
        assert result["file_type"] == "pptx"
        assert result["replacements"] == 1
        assert result["knowledge_synced"] is True
        assert result["source_sha256"] != read_result["source_sha256"]

        reopened = Presentation(deck_path)
        texts = [
            shape.text for slide in reopened.slides for shape in slide.shapes if getattr(shape, "has_text_frame", False)
        ]
        assert any("Client-ready narrative" in text for text in texts)
        assert (
            "Draft narrative" not in json.loads(await _read_file(entity_id=entity_id, path=deck_path.name))["content"]
        )
    finally:
        settings.MANOR_FS_ENABLED = old_enabled
        settings.MANOR_FS_ROOT = old_root


@pytest.mark.asyncio
async def test_read_file_rejects_stale_continuation_sha(tmp_path):
    settings = get_settings()
    old_enabled = settings.MANOR_FS_ENABLED
    old_root = settings.MANOR_FS_ROOT
    settings.MANOR_FS_ENABLED = True
    settings.MANOR_FS_ROOT = str(tmp_path)
    try:
        entity_id = "ent_1"
        entity_root = tmp_path / entity_id
        entity_root.mkdir()
        (entity_root / "note.txt").write_text("v1\n", encoding="utf-8")

        result = json.loads(await _read_file(entity_id=entity_id, path="note.txt"))
        (entity_root / "note.txt").write_text("v2\n", encoding="utf-8")

        stale = json.loads(
            await _read_file(
                entity_id=entity_id,
                path="note.txt",
                expected_sha256=result["source_sha256"],
            )
        )

        assert stale["error"] == "source_changed"
        assert stale["expected_sha256"] == result["source_sha256"]
        assert stale["source_sha256"] != result["source_sha256"]
    finally:
        settings.MANOR_FS_ENABLED = old_enabled
        settings.MANOR_FS_ROOT = old_root


@pytest.mark.asyncio
async def test_write_file_rejects_stale_expected_sha_without_overwrite(tmp_path):
    settings = get_settings()
    old_enabled = settings.MANOR_FS_ENABLED
    old_root = settings.MANOR_FS_ROOT
    settings.MANOR_FS_ENABLED = True
    settings.MANOR_FS_ROOT = str(tmp_path)
    try:
        entity_id = "ent_1"
        entity_root = tmp_path / entity_id
        entity_root.mkdir()
        target = entity_root / "draft.md"
        target.write_text("original\n", encoding="utf-8")
        read_result = json.loads(await _read_file(entity_id=entity_id, path="draft.md"))

        target.write_text("changed by someone else\n", encoding="utf-8")
        write_result = json.loads(
            await _write_file(
                entity_id=entity_id,
                path="draft.md",
                content="ai overwrite\n",
                expected_sha256=read_result["source_sha256"],
            )
        )

        assert write_result["error"] == "source_changed"
        assert target.read_text(encoding="utf-8") == "changed by someone else\n"
    finally:
        settings.MANOR_FS_ENABLED = old_enabled
        settings.MANOR_FS_ROOT = old_root


@pytest.mark.asyncio
async def test_write_file_scopes_new_workspace_files_to_workspace_folder(tmp_path, monkeypatch):
    settings = get_settings()
    old_enabled = settings.MANOR_FS_ENABLED
    old_root = settings.MANOR_FS_ROOT
    settings.MANOR_FS_ENABLED = True
    settings.MANOR_FS_ROOT = str(tmp_path)
    guard_call: dict[str, object] = {}
    sync_call: dict[str, object] = {}

    async def fake_guard_file_mutation(**kwargs):
        guard_call.update(kwargs)
        return None

    async def fake_sync_file_to_knowledge(**kwargs):
        sync_call.update(kwargs)
        return SimpleNamespace(synced=True, document_id="doc_1", reason="ok")

    async def fake_workspace_base_dir(**kwargs):
        assert kwargs["workspace_id"] == "ws_1"
        return "Workspaces/Launch Workspace"

    monkeypatch.setattr(file_tools_module, "runtime_guard_file_mutation", fake_guard_file_mutation)
    monkeypatch.setattr(file_tools_module, "runtime_sync_entity_file_to_knowledge", fake_sync_file_to_knowledge)
    monkeypatch.setattr(
        "packages.core.services.generated_media_naming.resolve_workspace_artifact_base_dir",
        fake_workspace_base_dir,
    )
    try:
        entity_id = "ent_1"
        entity_root = tmp_path / entity_id
        entity_root.mkdir()

        result = json.loads(
            await _write_file(
                entity_id=entity_id,
                path="brief.md",
                content="# Brief\n",
                workspace_id="ws_1",
                _user_id_from_context="user_1",
            )
        )

        expected_path = "Workspaces/Launch Workspace/documents/brief.md"
        assert result["written"] is True
        assert result["path"] == expected_path
        scoped_file = entity_root / "Workspaces" / "Launch Workspace" / "documents" / "brief.md"
        assert scoped_file.read_text(encoding="utf-8") == "# Brief\n"
        assert not (entity_root / "brief.md").exists()
        assert guard_call["paths"] == [expected_path]
        assert sync_call["workspace_id"] == "ws_1"
        assert sync_call["abs_path"] == str(scoped_file)
    finally:
        settings.MANOR_FS_ENABLED = old_enabled
        settings.MANOR_FS_ROOT = old_root


@pytest.mark.asyncio
async def test_write_file_keeps_existing_root_file_path_in_workspace(tmp_path, monkeypatch):
    settings = get_settings()
    old_enabled = settings.MANOR_FS_ENABLED
    old_root = settings.MANOR_FS_ROOT
    settings.MANOR_FS_ENABLED = True
    settings.MANOR_FS_ROOT = str(tmp_path)

    async def fake_guard_file_mutation(**_kwargs):
        return None

    async def fake_sync_file_to_knowledge(**_kwargs):
        return SimpleNamespace(synced=True, document_id="doc_1", reason="ok")

    async def fake_workspace_base_dir(**_kwargs):
        return "Workspaces/Launch Workspace"

    monkeypatch.setattr(file_tools_module, "runtime_guard_file_mutation", fake_guard_file_mutation)
    monkeypatch.setattr(file_tools_module, "runtime_sync_entity_file_to_knowledge", fake_sync_file_to_knowledge)
    monkeypatch.setattr(
        "packages.core.services.generated_media_naming.resolve_workspace_artifact_base_dir",
        fake_workspace_base_dir,
    )
    try:
        entity_id = "ent_1"
        entity_root = tmp_path / entity_id
        entity_root.mkdir()
        target = entity_root / "draft.md"
        target.write_text("old\n", encoding="utf-8")

        result = json.loads(
            await _write_file(
                entity_id=entity_id,
                path="draft.md",
                content="new\n",
                workspace_id="ws_1",
            )
        )

        assert result["written"] is True
        assert result["path"] == "draft.md"
        assert target.read_text(encoding="utf-8") == "new\n"
        assert not (entity_root / "Workspaces" / "Launch Workspace" / "documents" / "draft.md").exists()
    finally:
        settings.MANOR_FS_ENABLED = old_enabled
        settings.MANOR_FS_ROOT = old_root


@pytest.mark.asyncio
async def test_edit_file_uses_expected_sha_for_safe_mutation(tmp_path, monkeypatch):
    settings = get_settings()
    old_enabled = settings.MANOR_FS_ENABLED
    old_root = settings.MANOR_FS_ROOT
    settings.MANOR_FS_ENABLED = True
    settings.MANOR_FS_ROOT = str(tmp_path)
    monkeypatch.setattr(
        "packages.core.services.ai_file_permissions.guard_ai_file_mutation",
        AsyncMock(return_value=None),
    )
    try:
        entity_id = "ent_1"
        entity_root = tmp_path / entity_id
        entity_root.mkdir()
        target = entity_root / "draft.md"
        target.write_text("hello world\n", encoding="utf-8")
        read_result = json.loads(await _read_file(entity_id=entity_id, path="draft.md"))

        edit_result = json.loads(
            await _edit_file(
                entity_id=entity_id,
                path="draft.md",
                old_text="hello",
                new_text="hi",
                expected_sha256=read_result["source_sha256"],
            )
        )

        assert edit_result["edited"] is True
        assert target.read_text(encoding="utf-8") == "hi world\n"
        assert edit_result["source_sha256"] != read_result["source_sha256"]
    finally:
        settings.MANOR_FS_ENABLED = old_enabled
        settings.MANOR_FS_ROOT = old_root


@pytest.mark.asyncio
async def test_edit_file_rejects_stale_expected_sha_without_edit(tmp_path):
    settings = get_settings()
    old_enabled = settings.MANOR_FS_ENABLED
    old_root = settings.MANOR_FS_ROOT
    settings.MANOR_FS_ENABLED = True
    settings.MANOR_FS_ROOT = str(tmp_path)
    try:
        entity_id = "ent_1"
        entity_root = tmp_path / entity_id
        entity_root.mkdir()
        target = entity_root / "draft.md"
        target.write_text("hello world\n", encoding="utf-8")
        read_result = json.loads(await _read_file(entity_id=entity_id, path="draft.md"))

        target.write_text("hello changed\n", encoding="utf-8")
        edit_result = json.loads(
            await _edit_file(
                entity_id=entity_id,
                path="draft.md",
                old_text="hello",
                new_text="hi",
                expected_sha256=read_result["source_sha256"],
            )
        )

        assert edit_result["error"] == "source_changed"
        assert target.read_text(encoding="utf-8") == "hello changed\n"
    finally:
        settings.MANOR_FS_ENABLED = old_enabled
        settings.MANOR_FS_ROOT = old_root


@pytest.mark.asyncio
async def test_delete_file_rejects_stale_expected_sha_without_delete(tmp_path):
    settings = get_settings()
    old_enabled = settings.MANOR_FS_ENABLED
    old_root = settings.MANOR_FS_ROOT
    settings.MANOR_FS_ENABLED = True
    settings.MANOR_FS_ROOT = str(tmp_path)
    try:
        entity_id = "ent_1"
        entity_root = tmp_path / entity_id
        entity_root.mkdir()
        target = entity_root / "draft.md"
        target.write_text("delete me later\n", encoding="utf-8")
        read_result = json.loads(await _read_file(entity_id=entity_id, path="draft.md"))

        target.write_text("changed before delete\n", encoding="utf-8")
        delete_result = json.loads(
            await _delete_file(
                entity_id=entity_id,
                path="draft.md",
                expected_sha256=read_result["source_sha256"],
            )
        )

        assert delete_result["error"] == "source_changed"
        assert target.exists()
        assert target.read_text(encoding="utf-8") == "changed before delete\n"
    finally:
        settings.MANOR_FS_ENABLED = old_enabled
        settings.MANOR_FS_ROOT = old_root


@pytest.mark.asyncio
async def test_list_files_is_paginated_by_default(tmp_path):
    settings = get_settings()
    old_enabled = settings.MANOR_FS_ENABLED
    old_root = settings.MANOR_FS_ROOT
    settings.MANOR_FS_ENABLED = True
    settings.MANOR_FS_ROOT = str(tmp_path)
    try:
        entity_id = "ent_1"
        entity_root = tmp_path / entity_id
        entity_root.mkdir()
        for index in range(125):
            (entity_root / f"file_{index:03d}.txt").write_text("x", encoding="utf-8")

        first = json.loads(await _list_files(entity_id=entity_id))
        second = json.loads(await _list_files(entity_id=entity_id, offset=first["next_offset"]))

        assert first["count"] == 100
        assert first["has_more"] is True
        assert first["next_offset"] == 100
        assert first["total"] == 125
        assert second["count"] == 25
        assert second["has_more"] is False
        assert second["entries"][0]["path"] == "file_100.txt"
    finally:
        settings.MANOR_FS_ENABLED = old_enabled
        settings.MANOR_FS_ROOT = old_root


@pytest.mark.asyncio
async def test_glob_files_supports_pagination(tmp_path):
    settings = get_settings()
    old_enabled = settings.MANOR_FS_ENABLED
    old_root = settings.MANOR_FS_ROOT
    settings.MANOR_FS_ENABLED = True
    settings.MANOR_FS_ROOT = str(tmp_path)
    try:
        entity_id = "ent_1"
        entity_root = tmp_path / entity_id
        entity_root.mkdir()
        for index in range(25):
            (entity_root / f"note_{index:03d}.md").write_text("x", encoding="utf-8")

        first = json.loads(await _glob_files(entity_id=entity_id, pattern="*.md", limit=10))
        second = json.loads(
            await _glob_files(
                entity_id=entity_id,
                pattern="*.md",
                limit=10,
                offset=first["next_offset"],
            )
        )

        assert first["count"] == 10
        assert first["has_more"] is True
        assert first["next_offset"] == 10
        assert first["files"][0] == "note_000.md"
        assert second["count"] == 10
        assert second["files"][0] == "note_010.md"
    finally:
        settings.MANOR_FS_ENABLED = old_enabled
        settings.MANOR_FS_ROOT = old_root


@pytest.mark.asyncio
async def test_grep_files_supports_pagination(tmp_path):
    settings = get_settings()
    old_enabled = settings.MANOR_FS_ENABLED
    old_root = settings.MANOR_FS_ROOT
    settings.MANOR_FS_ENABLED = True
    settings.MANOR_FS_ROOT = str(tmp_path)
    try:
        entity_id = "ent_1"
        entity_root = tmp_path / entity_id
        entity_root.mkdir()
        for index in range(12):
            (entity_root / f"note_{index:03d}.txt").write_text(
                f"needle line {index}\n",
                encoding="utf-8",
            )

        first = json.loads(await _grep_files(entity_id=entity_id, pattern="needle", max_matches=5))
        second = json.loads(
            await _grep_files(
                entity_id=entity_id,
                pattern="needle",
                max_matches=5,
                offset=first["next_offset"],
            )
        )

        assert first["count"] == 5
        assert first["has_more"] is True
        assert first["next_offset"] == 5
        assert first["matches"][0]["file"] == "note_000.txt"
        assert second["count"] == 5
        assert second["matches"][0]["file"] == "note_005.txt"
    finally:
        settings.MANOR_FS_ENABLED = old_enabled
        settings.MANOR_FS_ROOT = old_root


@pytest.mark.asyncio
async def test_web_fetch_returns_digest_and_continuation(monkeypatch):
    html = "<html><body><h1>Title</h1><p>" + ("hello " * 5000) + "</p></body></html>"

    async def fake_fetch_url(url: str):
        return SimpleNamespace(
            content_type="text/html; charset=utf-8",
            content=html.encode("utf-8"),
        )

    monkeypatch.setattr("packages.core.services.web_fetch.fetch_url", fake_fetch_url)

    first = json.loads(
        await _web_fetch_handler(
            url="https://example.com/page",
            max_length=1200,
        )
    )
    second = json.loads(
        await _web_fetch_handler(
            url="https://example.com/page",
            max_length=1200,
            offset=first["next_offset"],
            expected_sha256=first["source_sha256"],
        )
    )

    assert first["url"] == "https://example.com/page"
    assert first["truncated"] is True
    assert first["next_offset"] == 1200
    assert len(first["source_sha256"]) == 64
    assert len(first["slice_sha256"]) == 64
    assert second["offset"] == first["next_offset"]
    assert second["source_sha256"] == first["source_sha256"]
    assert "expected_sha256=source_sha256" in first["hint"]


@pytest.mark.asyncio
async def test_web_fetch_rejects_stale_expected_sha(monkeypatch):
    calls = 0

    async def fake_fetch_url(url: str):
        nonlocal calls
        calls += 1
        text = "first version" if calls == 1 else "second version"
        return SimpleNamespace(
            content_type="text/plain",
            content=text.encode("utf-8"),
        )

    monkeypatch.setattr("packages.core.services.web_fetch.fetch_url", fake_fetch_url)

    first = json.loads(await _web_fetch_handler(url="https://example.com/page"))
    stale = json.loads(
        await _web_fetch_handler(
            url="https://example.com/page",
            expected_sha256=first["source_sha256"],
        )
    )

    assert stale["error"] == "source_changed"
    assert stale["expected_sha256"] == first["source_sha256"]
    assert stale["source_sha256"] != first["source_sha256"]


@pytest.mark.asyncio
async def test_web_fetch_flags_javascript_rendered_shell(monkeypatch):
    html = """<!doctype html>
    <html>
      <head><title>SPA</title><script type="module" src="/assets/app.js"></script></head>
      <body><div id="app"></div></body>
    </html>"""

    async def fake_fetch_url(url: str):
        return SimpleNamespace(
            content_type="text/html; charset=utf-8",
            content=html.encode("utf-8"),
        )

    monkeypatch.setattr("packages.core.services.web_fetch.fetch_url", fake_fetch_url)

    payload = json.loads(await _web_fetch_handler(url="https://example.com"))

    assert payload["url"] == "https://example.com"
    assert "dynamic_page_hint" in payload
    assert "browse_web" in payload["dynamic_page_hint"]


def test_document_tool_summary_omits_detail_fields_by_default():
    doc = SimpleNamespace(
        id="doc_1",
        name="brief.md",
        file_type="md",
        file_size=123,
        mime_type="text/markdown",
        source="upload",
        vector_status="ready",
        folder_id="folder_1",
        fs_path="notes/brief.md",
        created_at=None,
    )

    summary = _doc_to_dict(doc)
    details = _doc_to_dict(doc, detail="details")

    assert summary == {
        "id": "doc_1",
        "name": "brief.md",
        "file_type": "md",
        "file_size": 123,
    }
    assert details["mime_type"] == "text/markdown"
    assert details["vector_status"] == "ready"
    assert details["folder_id"] == "folder_1"
    assert details["fs_path"] == "notes/brief.md"


def test_manor_document_details_include_fs_path_only_when_requested():
    doc = SimpleNamespace(
        id="doc_1",
        name="brief.md",
        file_type="md",
        file_size=123,
        mime_type="text/markdown",
        source="upload",
        vector_status="ready",
        folder_id="folder_1",
        fs_path="notes/brief.md",
    )

    summary = _doc_summary(doc)
    details = _doc_summary(doc, details=True)

    assert "fs_path" not in summary
    assert details["fs_path"] == "notes/brief.md"


@pytest.mark.asyncio
async def test_generate_document_file_rejects_stale_expected_sha(tmp_path):
    settings = get_settings()
    old_enabled = settings.MANOR_FS_ENABLED
    old_root = settings.MANOR_FS_ROOT
    settings.MANOR_FS_ENABLED = True
    settings.MANOR_FS_ROOT = str(tmp_path)
    try:
        entity_id = "ent_1"
        entity_root = tmp_path / entity_id
        entity_root.mkdir()
        target = entity_root / "deliverable.md"
        target.write_text("v1\n", encoding="utf-8")
        read_result = json.loads(await _read_file(entity_id=entity_id, path="deliverable.md"))

        target.write_text("v2\n", encoding="utf-8")
        result = json.loads(
            await _generate_document_file(
                entity_id=entity_id,
                name="deliverable.md",
                content="new generated content\n",
                expected_sha256=read_result["source_sha256"],
            )
        )

        assert result["error"] == "source_changed"
        assert target.read_text(encoding="utf-8") == "v2\n"
    finally:
        settings.MANOR_FS_ENABLED = old_enabled
        settings.MANOR_FS_ROOT = old_root


def test_staff_summary_requires_details_for_extra_fields():
    staff = SimpleNamespace(
        id="staff_1",
        name="Simon",
        email="simon@example.com",
        kind="employee",
        title="PM",
        status="active",
        company_name="Manor",
        role_id="role_1",
        meta={"department": "Product", "role": "Lead"},
    )

    summary = _staff_summary(staff)
    details = _staff_summary(staff, details=True)

    assert summary == {
        "id": "staff_1",
        "name": "Simon",
        "email": "simon@example.com",
        "kind": "employee",
        "title": "PM",
        "status": "active",
    }
    assert details["company_name"] == "Manor"
    assert details["department"] == "Product"


def test_manor_schema_stays_compact_because_catalog_is_searchable():
    description = MANOR_SCHEMA["function"]["description"]

    assert len(description) < 900
    assert "action='search'" in description
    assert "create_scheduled_job" in description
    assert "list_files/glob_files/grep_files" not in description


def test_generate_file_schema_stays_compact_but_keeps_video_refs():
    schema_size = len(json.dumps(GENERATE_FILE_SCHEMA, ensure_ascii=False))
    props = GENERATE_FILE_SCHEMA["function"]["parameters"]["properties"]

    assert schema_size < 4000
    assert "code" in props["kind"]["enum"]
    assert props["duration"]["enum"] == props["params"]["properties"]["duration"]["enum"]
    assert "files" in props["params"]["properties"]
    assert "entry" in props["params"]["properties"]
    assert "first_frame_url" in props
    assert "last_frame_url" in props
    assert "reference_url" in props["params"]["properties"]
    assert "reference_urls" in props["params"]["properties"]
    assert "expected_sha256" in props


def test_search_tools_deferred_hint_stays_short():
    pool = ToolPool()
    pool._register_search_tools()
    for index in range(30):
        name = f"tool_{index:02d}"
        pool.register(name, _minimal_schema(name), handler=lambda: "ok")

    schemas, _ = _tool_surface_for_pool(pool, is_master=True)
    search_description = next(
        schema["function"]["description"] for schema in schemas if schema["function"]["name"] == "search_tools"
    )

    assert len(search_description) < 520
    assert "tool_11" in search_description
    assert "tool_12" not in search_description
    assert "(+18 more)" in search_description


def test_sandbox_save_result_schema_stays_compact():
    schema_size = len(json.dumps(_SANDBOX_SAVE_RESULT_SCHEMA, ensure_ascii=False))
    props = _SANDBOX_SAVE_RESULT_SCHEMA["function"]["parameters"]["properties"]

    assert schema_size < 900
    assert "filename" in _SANDBOX_SAVE_RESULT_SCHEMA["function"]["parameters"]["required"]
    assert {"sandbox_id", "file_path", "url", "filename"} <= set(props)


@pytest.mark.asyncio
async def test_bash_local_output_marks_truncation_with_digest(tmp_path):
    result = json.loads(
        await _execute_local(
            "python3 -c \"import sys; sys.stdout.write('x' * 70000)\"",
            timeout=5,
            cwd=str(tmp_path),
        )
    )

    assert result["exit_code"] == 0
    assert len(result["stdout"]) == 65536
    assert result["stdout_truncated"] is True
    assert result["stdout_chars"] == 70000
    assert len(result["stdout_sha256"]) == 64
    assert "read_file with offsets" in result["stdout_hint"]


@pytest.mark.asyncio
async def test_sandbox_read_file_returns_structured_metadata(monkeypatch):
    class FakeClient:
        async def read_file(self, sandbox_id: str, path: str):
            return SimpleNamespace(
                path=path,
                content="partial content",
                size=123456,
                truncated=True,
            )

        async def close(self):
            return None

    monkeypatch.setattr("packages.core.ai.tools.sandbox_tools._get_client", lambda: FakeClient())

    result = json.loads(await _sandbox_read_file(sandbox_id="sandbox_1", path="/skill/output.txt"))

    assert result["path"] == "/skill/output.txt"
    assert result["content"] == "partial content"
    assert result["size"] == 123456
    assert result["truncated"] is True
    assert len(result["content_sha256"]) == 64
    assert "truncated" in result["hint"]


def test_master_always_loaded_tool_schema_budget():
    pool = ToolPool()
    pool.initialize()
    schemas, _ = _tool_surface_for_pool(pool, is_master=True)
    schema_size = sum(len(json.dumps(schema, ensure_ascii=False)) for schema in schemas)
    names = {schema["function"]["name"] for schema in schemas}

    assert schema_size < 16_500
    assert "generate_file" in names
    assert "search_tools" in names


@pytest.mark.asyncio
async def test_tool_cache_keys_include_entity_version(monkeypatch):
    monkeypatch.setattr(
        "packages.core.services.tool_cache_version.get_tool_cache_version",
        AsyncMock(return_value=7),
    )

    doc_key = await document_cache_key("list_documents", "ent_1", {"limit": 20})
    manor_key = await _read_cache_key("ent_1", "list_staff", {"limit": 20})

    assert ":v7:" in doc_key
    assert ":v7:" in manor_key


@pytest.mark.asyncio
async def test_manor_cache_invalidation_bumps_shared_namespaces(monkeypatch):
    bumped: list[tuple[str, tuple[str, ...]]] = []

    async def fake_bump(entity_id: str, *namespaces: str) -> None:
        bumped.append((entity_id, namespaces))

    monkeypatch.setattr(
        "packages.core.services.tool_cache_version.bump_tool_cache_version",
        fake_bump,
    )

    await _invalidate_read_cache("ent_1", "list_documents", "search_documents", "list_staff")

    assert bumped == [("ent_1", ("documents", "staff"))]
