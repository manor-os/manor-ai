"""Text extraction from various file formats for RAG embedding.

Supports: plain text, markdown, HTML, PDF, CSV, JSON.
Each extractor returns the text content as a string.
"""
import csv
import asyncio
import io
import json
import logging
import os
import shutil
import subprocess
import tempfile
from pathlib import Path

logger = logging.getLogger(__name__)

# Maximum text length to extract (avoid huge embeddings)
MAX_EXTRACT_CHARS = 100_000
MAX_SPREADSHEET_EXTRACT_CHARS = 300_000
MAX_SPREADSHEET_ROWS_PER_SHEET = 5_000


async def extract_text(file_path: str, mime_type: str = None, file_type: str = None) -> str:
    """Extract text from a file based on its type.

    Priority: mime_type > file_type > file extension.
    Returns empty string if extraction fails or type is unsupported.
    """
    if not file_path or not os.path.isfile(file_path):
        return ""

    ext = Path(file_path).suffix.lower().lstrip(".")
    effective_type = file_type or ext

    try:
        if effective_type in ("txt", "text", "log", "cfg", "ini", "env", "sh", "bash", "zsh"):
            return _extract_plain_text(file_path)
        elif effective_type in ("md", "markdown"):
            return _extract_markdown(file_path)
        elif effective_type in ("html", "htm"):
            return _extract_html(file_path)
        elif effective_type in ("pdf",):
            return await _extract_pdf(file_path)
        elif effective_type in ("csv", "tsv"):
            return _extract_csv(file_path)
        elif effective_type in ("json", "jsonl"):
            return _extract_json(file_path)
        elif effective_type in ("py", "js", "ts", "tsx", "jsx", "java", "go", "rs", "c", "cpp", "h", "rb", "sql"):
            return _extract_plain_text(file_path)  # source code
        elif effective_type in ("docx", "wps"):
            return _extract_docx(file_path)
        elif effective_type in ("doc",):
            return await _extract_legacy_doc(file_path)
        elif effective_type in ("xlsx", "xls", "et"):
            return _extract_xlsx(file_path)
        elif effective_type in ("pptx", "ppt", "dps"):
            return _extract_pptx(file_path)
        elif effective_type in ("yaml", "yml", "toml"):
            return _extract_plain_text(file_path)  # config files
        elif mime_type and mime_type.startswith("text/"):
            return _extract_plain_text(file_path)
        else:
            logger.info("Unsupported file type for extraction: %s (mime=%s)", effective_type, mime_type)
            return ""
    except Exception as e:
        logger.error("Text extraction failed for %s: %s", file_path, e)
        return ""


def _extract_plain_text(path: str) -> str:
    """Read a plain text file."""
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()[:MAX_EXTRACT_CHARS]


def _extract_markdown(path: str) -> str:
    """Extract text from markdown (strip basic formatting)."""
    import re
    text = _extract_plain_text(path)
    # Strip common markdown syntax but keep content
    text = re.sub(r'#{1,6}\s+', '', text)  # headers
    text = re.sub(r'\*{1,2}([^*]+)\*{1,2}', r'\1', text)  # bold/italic
    text = re.sub(r'`{1,3}[^`]*`{1,3}', '', text)  # inline code
    text = re.sub(r'!\[.*?\]\(.*?\)', '', text)  # images
    text = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', text)  # links
    return text


def _extract_html(path: str) -> str:
    """Extract text from HTML (strip tags)."""
    import re
    raw = _extract_plain_text(path)
    # Remove script and style blocks
    raw = re.sub(r'<script[^>]*>.*?</script>', '', raw, flags=re.DOTALL | re.IGNORECASE)
    raw = re.sub(r'<style[^>]*>.*?</style>', '', raw, flags=re.DOTALL | re.IGNORECASE)
    # Strip all tags
    text = re.sub(r'<[^>]+>', ' ', raw)
    # Normalize whitespace
    text = re.sub(r'\s+', ' ', text).strip()
    return text[:MAX_EXTRACT_CHARS]


async def _extract_pdf(path: str) -> str:
    """Extract text from PDF using pypdf, PyPDF2, or pdfplumber.

    Tries libs in order of preference. ``pypdf`` is the maintained
    fork of PyPDF2. Returns empty string only when no supported PDF
    extraction backend is installed.
    """
    # Try pypdf (modern, maintained)
    try:
        from pypdf import PdfReader
        text_parts = []
        with open(path, "rb") as f:
            reader = PdfReader(f)
            for page in reader.pages[:200]:  # cap at 200 pages
                page_text = page.extract_text() or ""
                text_parts.append(page_text)
                if sum(len(t) for t in text_parts) > MAX_EXTRACT_CHARS:
                    break
        return "\n".join(text_parts)[:MAX_EXTRACT_CHARS]
    except ImportError:
        pass

    # Try PyPDF2 (deprecated but kept for legacy deploys)
    try:
        import PyPDF2
        text_parts = []
        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages[:200]:
                page_text = page.extract_text() or ""
                text_parts.append(page_text)
                if sum(len(t) for t in text_parts) > MAX_EXTRACT_CHARS:
                    break
        return "\n".join(text_parts)[:MAX_EXTRACT_CHARS]
    except ImportError:
        pass

    # Try pdfplumber
    try:
        import pdfplumber
        text_parts = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages[:200]:
                text_parts.append(page.extract_text() or "")
                if sum(len(t) for t in text_parts) > MAX_EXTRACT_CHARS:
                    break
        return "\n".join(text_parts)[:MAX_EXTRACT_CHARS]
    except ImportError:
        pass

    logger.warning("No PDF extraction library available. Install pypdf or pdfplumber.")
    return ""


def _extract_csv(path: str) -> str:
    """Convert CSV to readable text."""
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        rows = []
        for i, row in enumerate(reader):
            if i > 500:  # cap at 500 rows
                break
            rows.append(" | ".join(row))
    return "\n".join(rows)[:MAX_EXTRACT_CHARS]


def _extract_json(path: str) -> str:
    """Extract text from JSON — serialize as readable text."""
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        data = json.load(f)
    return json.dumps(data, indent=2, ensure_ascii=False)[:MAX_EXTRACT_CHARS]


def _extract_docx(path: str) -> str:
    """Extract text from .docx/.wps files using python-docx."""
    try:
        from docx import Document
        doc = Document(path)
        text_parts = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(text_parts)[:MAX_EXTRACT_CHARS]
    except ImportError:
        logger.warning("python-docx not installed. Install it for .docx/.wps extraction.")
        return ""


async def _extract_legacy_doc(path: str) -> str:
    """Extract text from legacy .doc files via LibreOffice conversion.

    ``python-docx`` cannot read binary Word 97-2003 ``.doc`` files. Keep this
    conversion in the API process so normal chat attachments do not need a
    sandbox just to read user-provided resumes/contracts.
    """
    return await asyncio.to_thread(_extract_legacy_doc_sync, path)


def _extract_legacy_doc_sync(path: str) -> str:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        logger.warning("LibreOffice/soffice not installed. Cannot extract legacy .doc files.")
        return ""

    with tempfile.TemporaryDirectory() as tmpdir:
        try:
            subprocess.run(
                [soffice, "--headless", "--convert-to", "docx", "--outdir", tmpdir, path],
                check=False,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                timeout=60,
            )
        except Exception as exc:
            logger.warning("Legacy .doc conversion to docx failed for %s: %s", path, exc)

        converted_docx = os.path.join(tmpdir, f"{Path(path).stem}.docx")
        if os.path.isfile(converted_docx):
            text = _extract_docx(converted_docx)
            if text.strip():
                return text[:MAX_EXTRACT_CHARS]

        try:
            subprocess.run(
                [soffice, "--headless", "--convert-to", "txt:Text", "--outdir", tmpdir, path],
                check=False,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                timeout=60,
            )
        except Exception as exc:
            logger.warning("Legacy .doc conversion to text failed for %s: %s", path, exc)
            return ""

        converted_txt = os.path.join(tmpdir, f"{Path(path).stem}.txt")
        if os.path.isfile(converted_txt):
            return _extract_plain_text(converted_txt)[:MAX_EXTRACT_CHARS]

    logger.warning("Legacy .doc conversion produced no readable output for %s", path)
    return ""


def _extract_xlsx(path: str) -> str:
    """Extract text from .xlsx/.et files using openpyxl."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(path, read_only=True, data_only=True)
        text_parts = []
        for ws in wb.worksheets[:20]:  # cap at 20 sheets
            text_parts.append(f"[Sheet: {ws.title}]")
            for row in ws.iter_rows(max_row=MAX_SPREADSHEET_ROWS_PER_SHEET, values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    text_parts.append(" | ".join(cells))
            if sum(len(t) for t in text_parts) > MAX_SPREADSHEET_EXTRACT_CHARS:
                break
        wb.close()
        return "\n".join(text_parts)[:MAX_SPREADSHEET_EXTRACT_CHARS]
    except ImportError:
        logger.warning("openpyxl not installed. Install it for .xlsx/.et extraction.")
        return ""


def _extract_pptx(path: str) -> str:
    """Extract text from .pptx/.dps files using python-pptx."""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        text_parts = []
        for i, slide in enumerate(prs.slides, 1):  # cap at 100 slides
            if i > 100:
                break
            text_parts.append(f"[Slide {i}]")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            text_parts.append(para.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            text_parts.append(" | ".join(cells))
            if sum(len(t) for t in text_parts) > MAX_EXTRACT_CHARS:
                break
        return "\n".join(text_parts)[:MAX_EXTRACT_CHARS]
    except ImportError:
        logger.warning("python-pptx not installed. Install it for .pptx/.dps extraction.")
        return ""
