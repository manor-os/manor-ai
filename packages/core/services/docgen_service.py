"""Document generation service — creates Word, PDF, and PowerPoint files from structured content."""
from __future__ import annotations

import io
import logging
import re
from typing import Any

from packages.core.ai.runtime import runtime_execute_docgen_completion

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# MIME types and format metadata
# ---------------------------------------------------------------------------

FORMAT_META: dict[str, dict[str, str]] = {
    "docx": {
        "name": "Word Document",
        "mime": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "ext": ".docx",
    },
    "pdf": {
        "name": "PDF Document",
        "mime": "application/pdf",
        "ext": ".pdf",
    },
    "pptx": {
        "name": "PowerPoint Presentation",
        "mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "ext": ".pptx",
    },
}

SUPPORTED_FORMATS = list(FORMAT_META.keys())


def _check_format(fmt: str) -> None:
    if fmt not in FORMAT_META:
        raise ValueError(f"Unsupported format '{fmt}'. Choose from: {', '.join(SUPPORTED_FORMATS)}")


def _slugify(title: str) -> str:
    """Turn a title into a safe filename stem."""
    slug = re.sub(r"[^\w\s-]", "", title.strip().lower())
    slug = re.sub(r"[\s]+", "_", slug)
    return slug[:80] or "document"


# ---------------------------------------------------------------------------
# Markdown parsing helpers
# ---------------------------------------------------------------------------

_HEADING_RE = re.compile(r"^(#{1,6})\s+(.*)", re.MULTILINE)
_BULLET_RE = re.compile(r"^\s*[-*+]\s+(.*)")
_NUMBERED_RE = re.compile(r"^\s*\d+\.\s+(.*)")
_TABLE_SEP_RE = re.compile(r"^\s*\|?\s*[-:]+[-| :]*$")
_TABLE_ROW_RE = re.compile(r"^\s*\|(.+)\|\s*$")
_BOLD_RE = re.compile(r"\*\*(.+?)\*\*")
_ITALIC_RE = re.compile(r"\*(.+?)\*")
_UNDERLINE_RE = re.compile(r"\[\[U\]\](.+?)\[\[/U\]\]")
_STRIKE_RE = re.compile(r"\[\[S\]\](.+?)\[\[/S\]\]")


def _parse_blocks(content: str) -> list[dict[str, Any]]:
    """Parse markdown-ish content into a list of block dicts.

    Block types: heading, paragraph, bullet, numbered, table.
    """
    blocks: list[dict[str, Any]] = []
    lines = content.split("\n")
    i = 0
    while i < len(lines):
        line = lines[i]
        stripped_line = line.strip()

        if stripped_line == "[[PAGE_BREAK]]":
            blocks.append({"type": "page_break"})
            i += 1
            continue

        # Heading
        m = _HEADING_RE.match(line)
        if m:
            blocks.append({"type": "heading", "level": len(m.group(1)), "text": m.group(2).strip()})
            i += 1
            continue

        # Table — collect rows
        if _TABLE_ROW_RE.match(line):
            rows: list[list[str]] = []
            while i < len(lines):
                row_m = _TABLE_ROW_RE.match(lines[i])
                if row_m:
                    if not _TABLE_SEP_RE.match(lines[i]):
                        cells = [c.strip() for c in row_m.group(1).split("|")]
                        rows.append(cells)
                    i += 1
                else:
                    break
            if rows:
                blocks.append({"type": "table", "rows": rows})
            continue

        # Bullet
        bm = _BULLET_RE.match(line)
        if bm:
            items: list[str] = []
            while i < len(lines) and _BULLET_RE.match(lines[i]):
                items.append(_BULLET_RE.match(lines[i]).group(1))  # type: ignore[union-attr]
                i += 1
            blocks.append({"type": "bullet", "items": items})
            continue

        # Numbered list
        nm = _NUMBERED_RE.match(line)
        if nm:
            items_n: list[str] = []
            while i < len(lines) and _NUMBERED_RE.match(lines[i]):
                items_n.append(_NUMBERED_RE.match(lines[i]).group(1))  # type: ignore[union-attr]
                i += 1
            blocks.append({"type": "numbered", "items": items_n})
            continue

        # Blank line — skip
        if not line.strip():
            i += 1
            continue

        # Plain paragraph — collect consecutive non-blank non-special lines
        para_lines: list[str] = []
        while i < len(lines):
            ln = lines[i]
            if (
                not ln.strip()
                or _HEADING_RE.match(ln)
                or _BULLET_RE.match(ln)
                or _NUMBERED_RE.match(ln)
                or _TABLE_ROW_RE.match(ln)
            ):
                break
            para_lines.append(ln)
            i += 1
        if para_lines:
            blocks.append({"type": "paragraph", "text": " ".join(para_lines)})

    return blocks


def _strip_md_inline(text: str) -> str:
    """Remove bold/italic markdown for plain-text contexts."""
    text = _UNDERLINE_RE.sub(r"\1", text)
    text = _STRIKE_RE.sub(r"\1", text)
    text = _BOLD_RE.sub(r"\1", text)
    text = _ITALIC_RE.sub(r"\1", text)
    return text


# ---------------------------------------------------------------------------
# DOCX generation
# ---------------------------------------------------------------------------

async def generate_docx(title: str, content: str, options: dict | None = None) -> bytes:
    """Generate a Word document from markdown/text content.

    Requires ``python-docx``.
    """
    try:
        from docx import Document as DocxDocument  # type: ignore[import-untyped]
        from docx.shared import Pt  # type: ignore[import-untyped]
        from docx.enum.text import WD_ALIGN_PARAGRAPH  # type: ignore[import-untyped]
    except ImportError:
        raise RuntimeError(
            "python-docx is required for DOCX generation. Install it: pip install python-docx"
        )

    opts = options or {}
    doc = DocxDocument()

    # Title
    title_para = doc.add_heading(title, level=0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    blocks = _parse_blocks(content)

    for block in blocks:
        btype = block["type"]

        if btype == "heading":
            level = min(block["level"], 4)  # docx supports heading 1-9
            doc.add_heading(block["text"], level=level)

        elif btype == "paragraph":
            p = doc.add_paragraph()
            _apply_inline_formatting(p, block["text"])

        elif btype == "bullet":
            for item in block["items"]:
                p = doc.add_paragraph(style="List Bullet")
                _apply_inline_formatting(p, item)

        elif btype == "numbered":
            for item in block["items"]:
                p = doc.add_paragraph(style="List Number")
                _apply_inline_formatting(p, item)

        elif btype == "table":
            rows = block["rows"]
            if not rows:
                continue
            n_cols = max(len(r) for r in rows)
            table = doc.add_table(rows=len(rows), cols=n_cols)
            table.style = "Table Grid"
            for ri, row in enumerate(rows):
                for ci, cell_text in enumerate(row):
                    if ci < n_cols:
                        table.rows[ri].cells[ci].text = _strip_md_inline(cell_text)

        elif btype == "page_break":
            doc.add_page_break()

    # Apply default font size if requested
    font_size = opts.get("font_size", 11)
    for para in doc.paragraphs:
        for run in para.runs:
            run.font.size = Pt(font_size)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _apply_inline_formatting(paragraph, text: str) -> None:
    """Add runs to a paragraph with bold/italic from markdown inline syntax."""
    try:
        import docx  # noqa: F401  # type: ignore[import-untyped]
    except ImportError:
        paragraph.add_run(text)
        return

    bold = False
    italic = False
    underline = False
    strike = False
    i = 0
    markers = ("**", "*", "[[U]]", "[[/U]]", "[[S]]", "[[/S]]")

    def add_run(value: str) -> None:
        if not value:
            return
        run = paragraph.add_run(value)
        run.bold = bold
        run.italic = italic
        run.underline = underline
        run.font.strike = strike

    while i < len(text):
        if text.startswith("**", i):
            bold = not bold
            i += 2
            continue
        if text.startswith("*", i):
            italic = not italic
            i += 1
            continue
        if text.startswith("[[U]]", i):
            underline = True
            i += 5
            continue
        if text.startswith("[[/U]]", i):
            underline = False
            i += 6
            continue
        if text.startswith("[[S]]", i):
            strike = True
            i += 5
            continue
        if text.startswith("[[/S]]", i):
            strike = False
            i += 6
            continue

        next_positions = [text.find(marker, i + 1) for marker in markers]
        next_positions = [pos for pos in next_positions if pos >= 0]
        next_i = min(next_positions) if next_positions else len(text)
        add_run(text[i:next_i])
        i = next_i


# ---------------------------------------------------------------------------
# PDF generation
# ---------------------------------------------------------------------------

async def generate_pdf(title: str, content: str, options: dict | None = None) -> bytes:
    """Generate a PDF from markdown/text content.

    Tries reportlab first, then falls back to a basic HTML-wrapped approach.
    """
    # Try reportlab
    try:
        return _generate_pdf_reportlab(title, content, options)
    except ImportError:
        pass

    # Fallback: generate a simple PDF-like representation using FPDF2
    try:
        return _generate_pdf_fpdf2(title, content, options)
    except ImportError:
        pass

    # Last resort: return HTML bytes with a note
    logger.warning("No PDF library available. Install reportlab or fpdf2: pip install reportlab")
    raise RuntimeError(
        "No PDF generation library available. Install reportlab (pip install reportlab) "
        "or fpdf2 (pip install fpdf2)."
    )


def _generate_pdf_reportlab(title: str, content: str, options: dict | None = None) -> bytes:
    """Generate PDF using reportlab."""
    from reportlab.lib.pagesizes import letter  # type: ignore[import-untyped]
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle  # type: ignore[import-untyped]
    from reportlab.lib.units import inch  # type: ignore[import-untyped]
    from reportlab.platypus import (  # type: ignore[import-untyped]
        SimpleDocTemplate, Paragraph, Spacer, Table as RLTable, TableStyle,
        ListFlowable, ListItem,
    )
    from reportlab.lib import colors  # type: ignore[import-untyped]

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=letter, topMargin=0.75 * inch, bottomMargin=0.75 * inch)
    styles = getSampleStyleSheet()

    # Custom styles
    title_style = ParagraphStyle(
        "DocTitle", parent=styles["Title"], fontSize=20, spaceAfter=20,
    )
    h2_style = ParagraphStyle("H2", parent=styles["Heading2"], fontSize=16, spaceAfter=10)
    h3_style = ParagraphStyle("H3", parent=styles["Heading3"], fontSize=13, spaceAfter=8)
    body_style = styles["BodyText"]
    bullet_style = styles["BodyText"]

    heading_styles = {1: title_style, 2: h2_style, 3: h3_style}

    story: list = []
    story.append(Paragraph(title, title_style))
    story.append(Spacer(1, 12))

    blocks = _parse_blocks(content)

    for block in blocks:
        btype = block["type"]

        if btype == "heading":
            level = block["level"]
            style = heading_styles.get(level, h3_style)
            story.append(Paragraph(_strip_md_inline(block["text"]), style))
            story.append(Spacer(1, 6))

        elif btype == "paragraph":
            text = _md_to_rl_markup(block["text"])
            story.append(Paragraph(text, body_style))
            story.append(Spacer(1, 6))

        elif btype == "bullet":
            items = [
                ListItem(Paragraph(_md_to_rl_markup(item), bullet_style))
                for item in block["items"]
            ]
            story.append(ListFlowable(items, bulletType="bullet"))
            story.append(Spacer(1, 6))

        elif btype == "numbered":
            items = [
                ListItem(Paragraph(_md_to_rl_markup(item), bullet_style))
                for item in block["items"]
            ]
            story.append(ListFlowable(items, bulletType="1"))
            story.append(Spacer(1, 6))

        elif btype == "table":
            rows = block["rows"]
            if not rows:
                continue
            table_data = [[_strip_md_inline(c) for c in row] for row in rows]
            t = RLTable(table_data)
            t.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#4472C4")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#D9E2F3")]),
            ]))
            story.append(t)
            story.append(Spacer(1, 8))

    doc.build(story)
    return buf.getvalue()


def _md_to_rl_markup(text: str) -> str:
    """Convert markdown bold/italic to reportlab XML markup."""
    text = _BOLD_RE.sub(r"<b>\1</b>", text)
    text = _ITALIC_RE.sub(r"<i>\1</i>", text)
    # Escape angle brackets that aren't our tags
    return text


def _generate_pdf_fpdf2(title: str, content: str, options: dict | None = None) -> bytes:
    """Generate PDF using fpdf2."""
    from fpdf import FPDF  # type: ignore[import-untyped]

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    # Title
    pdf.set_font("Helvetica", "B", 18)
    pdf.cell(0, 12, title, ln=True, align="C")
    pdf.ln(8)

    blocks = _parse_blocks(content)

    for block in blocks:
        btype = block["type"]

        if btype == "heading":
            size = {1: 18, 2: 16, 3: 14}.get(block["level"], 12)
            pdf.set_font("Helvetica", "B", size)
            pdf.cell(0, 10, _strip_md_inline(block["text"]), ln=True)
            pdf.ln(3)

        elif btype == "paragraph":
            pdf.set_font("Helvetica", "", 11)
            pdf.multi_cell(0, 6, _strip_md_inline(block["text"]))
            pdf.ln(3)

        elif btype in ("bullet", "numbered"):
            pdf.set_font("Helvetica", "", 11)
            for idx, item in enumerate(block["items"]):
                prefix = f"  {idx + 1}. " if btype == "numbered" else "  - "
                pdf.multi_cell(0, 6, prefix + _strip_md_inline(item))
            pdf.ln(3)

        elif btype == "table":
            rows = block["rows"]
            if not rows:
                continue
            n_cols = max(len(r) for r in rows)
            col_w = (pdf.w - 20) / max(n_cols, 1)
            for ri, row in enumerate(rows):
                pdf.set_font("Helvetica", "B" if ri == 0 else "", 10)
                for ci in range(n_cols):
                    cell_text = _strip_md_inline(row[ci]) if ci < len(row) else ""
                    pdf.cell(col_w, 7, cell_text, border=1)
                pdf.ln()
            pdf.ln(3)

    return bytes(pdf.output())


# ---------------------------------------------------------------------------
# PPTX generation
# ---------------------------------------------------------------------------

async def generate_pptx(title: str, content: str, options: dict | None = None) -> bytes:
    """Generate a styled PowerPoint presentation from structured content.

    Requires ``python-pptx``.
    Splits content by ## headings into slides with backgrounds, accent bars,
    and formatted text.
    """
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
        from pptx.util import Inches, Pt, Emu  # type: ignore[import-untyped]
        from pptx.dml.color import RGBColor  # type: ignore[import-untyped]
        from pptx.enum.text import PP_ALIGN  # type: ignore[import-untyped]
    except ImportError:
        raise RuntimeError(
            "python-pptx is required for PPTX generation. Install it: pip install python-pptx"
        )

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # Accent palette
    ACCENTS = [
        RGBColor(0x0D, 0x94, 0x88),  # teal
        RGBColor(0x25, 0x63, 0xEB),  # blue
        RGBColor(0x7C, 0x3A, 0xED),  # violet
        RGBColor(0xDC, 0x26, 0x26),  # red
        RGBColor(0xD9, 0x77, 0x06),  # amber
        RGBColor(0x05, 0x96, 0x69),  # emerald
    ]

    blank_layout = prs.slide_layouts[6]  # Blank layout

    def _add_bg_rect(slide, color: RGBColor):
        """Add a full-slide background rectangle."""
        from pptx.util import Emu as E  # type: ignore[import-untyped]
        bg = slide.shapes.add_shape(1, E(0), E(0), slide_w, slide_h)  # 1 = rectangle
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()

    def _add_accent_bar(slide, accent: RGBColor, y=0, h=Inches(0.12)):
        """Add a thin accent bar across the top of the slide."""
        bar = slide.shapes.add_shape(1, Emu(0), Emu(y), slide_w, h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

    def _set_text(tf, text: str, size: int = 18, color: RGBColor | None = None,
                  bold: bool = False, align=None, font_name: str = "Calibri"):
        """Set text in a text frame with formatting."""
        tf.clear()
        p = tf.paragraphs[0]
        p.text = _strip_md_inline(text)
        p.font.size = Pt(size)
        p.font.name = font_name
        if color:
            p.font.color.rgb = color
        p.font.bold = bold
        if align:
            p.alignment = align

    # ── Title slide ──
    slide = prs.slides.add_slide(blank_layout)
    _add_bg_rect(slide, RGBColor(0x0F, 0x17, 0x2A))  # dark navy

    # Gradient overlay (accent bar at bottom)
    accent_bar = slide.shapes.add_shape(1, Emu(0), int(slide_h * 0.85), slide_w, int(slide_h * 0.15))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENTS[0]
    accent_bar.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(
        int(slide_w * 0.1), int(slide_h * 0.3),
        int(slide_w * 0.8), int(slide_h * 0.3),
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    _set_text(tf, title, size=44, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True,
              align=PP_ALIGN.CENTER, font_name="Calibri Light")

    # Subtitle if provided
    subtitle_text = options.get("subtitle", "") if options else ""
    if subtitle_text:
        sub_box = slide.shapes.add_textbox(
            int(slide_w * 0.15), int(slide_h * 0.6),
            int(slide_w * 0.7), int(slide_h * 0.1),
        )
        _set_text(sub_box.text_frame, subtitle_text, size=20,
                  color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.CENTER)

    # ── Content slides ──
    slides_data = _split_into_slides(content)

    for idx, slide_info in enumerate(slides_data):
        accent = ACCENTS[idx % len(ACCENTS)]
        slide = prs.slides.add_slide(blank_layout)

        # White background
        _add_bg_rect(slide, RGBColor(0xFF, 0xFF, 0xFF))

        # Top accent bar
        _add_accent_bar(slide, accent)

        # Slide title
        title_box = slide.shapes.add_textbox(
            int(slide_w * 0.06), int(slide_h * 0.05),
            int(slide_w * 0.88), int(slide_h * 0.12),
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        _set_text(tf, slide_info["title"], size=32, color=RGBColor(0x0F, 0x17, 0x2A),
                  bold=True, font_name="Calibri Light")

        # Left accent strip
        strip = slide.shapes.add_shape(
            1, int(slide_w * 0.06), int(slide_h * 0.18),
            Inches(0.06), int(slide_h * 0.72),
        )
        strip.fill.solid()
        strip.fill.fore_color.rgb = accent
        strip.line.fill.background()

        # Content area
        content_box = slide.shapes.add_textbox(
            int(slide_w * 0.09), int(slide_h * 0.20),
            int(slide_w * 0.85), int(slide_h * 0.72),
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.auto_size = None

        first = True
        for item in slide_info["items"]:
            cleaned = _strip_md_inline(item)
            if not cleaned:
                continue

            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()

            # Detect bullet items
            is_bullet = item.lstrip().startswith(("- ", "* ", "• ", "· "))
            is_numbered = bool(re.match(r"^\d+\.\s", item.lstrip()))

            if is_bullet:
                cleaned = re.sub(r"^[-*•·]\s+", "", cleaned)
                p.text = cleaned
                p.level = 1
            elif is_numbered:
                p.text = cleaned
                p.level = 1
            else:
                p.text = cleaned
                p.level = 0

            p.font.size = Pt(18)
            p.font.name = "Calibri"
            p.font.color.rgb = RGBColor(0x33, 0x41, 0x55)  # slate-700
            p.space_after = Pt(6)

            # Bold sub-headings (lines starting with emoji or all-caps-ish)
            if cleaned and (cleaned[0] in "🎨🧠⚡💡📌🔑🎯📋🎭⚙️🔧🔊🎙️📐📱🇨🇳🌏🎮🌐📖🚀💎✅🔥💬📊🤝🎬🎵🔉"
                           or cleaned.startswith("##")):
                p.font.bold = True
                p.font.color.rgb = RGBColor(0x1E, 0x29, 0x3B)  # slate-800
                p.font.size = Pt(20)

        # Slide number
        num_box = slide.shapes.add_textbox(
            int(slide_w * 0.92), int(slide_h * 0.93),
            int(slide_w * 0.06), int(slide_h * 0.05),
        )
        _set_text(num_box.text_frame, str(idx + 1), size=11,
                  color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.RIGHT)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _split_into_slides(content: str) -> list[dict[str, Any]]:
    """Split content by ## headings or --- Slide N --- markers into slide dicts."""
    # Detect "--- Slide N ---" format
    slide_marker = re.compile(r"---\s*Slide\s+\d+\s*---")
    if slide_marker.search(content):
        return _split_slide_markers(content, slide_marker)

    slides: list[dict[str, Any]] = []
    current: dict[str, Any] | None = None

    for line in content.split("\n"):
        m = re.match(r"^##\s+(.*)", line)
        if m:
            if current:
                slides.append(current)
            current = {"title": m.group(1).strip(), "items": []}
            continue

        # Skip top-level heading (used as presentation title)
        if re.match(r"^#\s+", line):
            continue

        if current is None:
            # Content before first ## heading — create an "Overview" slide
            if line.strip():
                current = {"title": "Overview", "items": []}
                current["items"].append(line.strip())
            continue

        # Collect non-empty lines as bullet items
        stripped = line.strip()
        if stripped:
            # Remove bullet/number prefix
            stripped = re.sub(r"^[-*+]\s+", "", stripped)
            stripped = re.sub(r"^\d+\.\s+", "", stripped)
            current["items"].append(stripped)

    if current:
        slides.append(current)

    return slides


def _split_slide_markers(content: str, marker: re.Pattern) -> list[dict[str, Any]]:
    """Split content by --- Slide N --- markers into slide dicts."""
    blocks = marker.split(content)
    slides: list[dict[str, Any]] = []

    for block in blocks:
        lines = [line.strip() for line in block.strip().split("\n") if line.strip()]
        if not lines:
            continue

        # First ## heading or first non-empty line becomes title
        title = "Slide"
        items: list[str] = []
        for line in lines:
            m = re.match(r"^##\s+(.*)", line)
            if m and title == "Slide":
                title = m.group(1).strip()
            elif re.match(r"^#\s+", line):
                continue  # skip top-level headings
            else:
                cleaned = re.sub(r"^[-*+]\s+", "", line)
                cleaned = re.sub(r"^\d+\.\s+", "", cleaned)
                items.append(cleaned)

        slides.append({"title": title, "items": items})

    return slides


# ---------------------------------------------------------------------------
# HTML preview
# ---------------------------------------------------------------------------

def content_to_html(title: str, content: str) -> str:
    """Render markdown-ish content to a styled HTML string for preview."""
    return _content_to_html(title, content)


def _content_to_html(title: str, content: str) -> str:
    """Internal HTML renderer."""
    blocks = _parse_blocks(content)
    parts = [
        "<!DOCTYPE html>",
        '<html><head><meta charset="utf-8">',
        "<style>",
        "body{font-family:system-ui,-apple-system,sans-serif;max-width:800px;margin:2rem auto;padding:0 1rem;color:#1a1a1a;line-height:1.6}",
        "h1{border-bottom:2px solid #4472C4;padding-bottom:0.3rem}",
        "h2{color:#2c5282}h3{color:#4a5568}",
        "table{border-collapse:collapse;width:100%;margin:1rem 0}",
        "th,td{border:1px solid #cbd5e0;padding:0.5rem;text-align:left}",
        "th{background:#4472C4;color:white}",
        "tr:nth-child(even){background:#f7fafc}",
        "ul,ol{padding-left:1.5rem}",
        "</style></head><body>",
        f"<h1>{_html_escape(title)}</h1>",
    ]

    for block in blocks:
        btype = block["type"]
        if btype == "heading":
            tag = f"h{min(block['level'] + 1, 6)}"
            parts.append(f"<{tag}>{_html_escape(block['text'])}</{tag}>")
        elif btype == "paragraph":
            parts.append(f"<p>{_md_to_html_inline(block['text'])}</p>")
        elif btype == "bullet":
            parts.append("<ul>")
            for item in block["items"]:
                parts.append(f"<li>{_md_to_html_inline(item)}</li>")
            parts.append("</ul>")
        elif btype == "numbered":
            parts.append("<ol>")
            for item in block["items"]:
                parts.append(f"<li>{_md_to_html_inline(item)}</li>")
            parts.append("</ol>")
        elif btype == "table":
            rows = block["rows"]
            parts.append("<table>")
            for ri, row in enumerate(rows):
                tag = "th" if ri == 0 else "td"
                parts.append("<tr>")
                for cell in row:
                    parts.append(f"<{tag}>{_html_escape(cell)}</{tag}>")
                parts.append("</tr>")
            parts.append("</table>")

    parts.append("</body></html>")
    return "\n".join(parts)


def _html_escape(text: str) -> str:
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def _md_to_html_inline(text: str) -> str:
    text = _html_escape(text)
    text = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", text)
    text = re.sub(r"\*(.+?)\*", r"<em>\1</em>", text)
    return text


# ---------------------------------------------------------------------------
# Availability check
# ---------------------------------------------------------------------------

def get_available_formats() -> list[dict[str, Any]]:
    """Return format metadata with availability flags."""
    results = []
    for fmt, meta in FORMAT_META.items():
        available = True
        reason = None
        if fmt == "docx":
            try:
                import docx  # type: ignore[import-untyped]  # noqa: F401
            except ImportError:
                available = False
                reason = "pip install python-docx"
        elif fmt == "pdf":
            has_rl = has_fpdf = False
            try:
                import reportlab  # type: ignore[import-untyped]  # noqa: F401
                has_rl = True
            except ImportError:
                pass
            try:
                import fpdf  # type: ignore[import-untyped]  # noqa: F401
                has_fpdf = True
            except ImportError:
                pass
            if not has_rl and not has_fpdf:
                available = False
                reason = "pip install reportlab (or fpdf2)"
        elif fmt == "pptx":
            try:
                import pptx  # type: ignore[import-untyped]  # noqa: F401
            except ImportError:
                available = False
                reason = "pip install python-pptx"
        results.append({
            "format": fmt,
            "name": meta["name"],
            "mime_type": meta["mime"],
            "extension": meta["ext"],
            "available": available,
            **({"install": reason} if reason else {}),
        })
    return results


# ---------------------------------------------------------------------------
# Main entry points
# ---------------------------------------------------------------------------

async def generate_document(
    entity_id: str,
    user_id: str,
    title: str,
    content: str,
    format: str,
    template: str | None = None,
    options: dict | None = None,
) -> tuple[bytes, str, str]:
    """Generate a document file.

    Returns ``(file_bytes, filename, mime_type)``.
    """
    _check_format(format)
    meta = FORMAT_META[format]

    if format == "docx":
        file_bytes = await generate_docx(title, content, options)
    elif format == "pdf":
        file_bytes = await generate_pdf(title, content, options)
    elif format == "pptx":
        file_bytes = await generate_pptx(title, content, options)
    else:
        raise ValueError(f"Unsupported format: {format}")

    filename = f"{_slugify(title)}{meta['ext']}"
    return file_bytes, filename, meta["mime"]


async def ai_generate_document(
    entity_id: str,
    user_id: str,
    prompt: str,
    format: str,
    options: dict | None = None,
) -> tuple[bytes, str, str]:
    """Use AI to generate document content, then render to the requested format.

    Calls the LLM to produce structured markdown content based on the prompt,
    then passes it to :func:`generate_document`.
    """
    _check_format(format)

    completion = await runtime_execute_docgen_completion(
        entity_id=entity_id,
        prompt=prompt,
        format_name=format,
    )

    content = completion.content

    # Extract title from first heading or use a default
    title_match = re.match(r"^#\s+(.+)", content, re.MULTILINE)
    title = title_match.group(1).strip() if title_match else "Generated Document"

    return await generate_document(
        entity_id=entity_id,
        user_id=user_id,
        title=title,
        content=content,
        format=format,
        options=options,
    )
