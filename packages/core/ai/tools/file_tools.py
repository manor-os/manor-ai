"""File operation tools — read, write, list, glob, grep.

All operations are scoped to the entity's directory: {MANOR_FS_ROOT}/{entity_id}/
"""
from __future__ import annotations

import asyncio
import base64
import fnmatch
import hashlib
import io
import json
import logging
import os
import re
import time
from typing import Any

from packages.core.ai.runtime.file_actions import (
    runtime_entity_file_root,
    runtime_guard_file_mutation,
    runtime_normalize_entity_file_path,
    runtime_sync_entity_file_to_knowledge,
    runtime_trash_knowledge_path,
    runtime_user_visible_file_path,
    runtime_write_entity_file_atomic,
)
from packages.core.ai.runtime.tool_context import runtime_tool_call_context_from_kwargs
from packages.core.services.workspace_layout import WorkspaceArtifactDir

logger = logging.getLogger(__name__)

READ_FILE_DEFAULT_LINES = 120
READ_FILE_MAX_LINES = 1000
READ_FILE_DEFAULT_CHARS = 12_000
READ_FILE_MAX_CHARS = 40_000
READ_FILE_MAX_IMAGE_BYTES = 12_000_000
LIST_FILES_DEFAULT_LIMIT = 100
LIST_FILES_MAX_LIMIT = 200
GLOB_FILES_DEFAULT_LIMIT = 100
GREP_FILES_DEFAULT_LIMIT = 50

# grep_files scan bounds. Production incident: a 1.9 GB / 2410-file workspace
# was 99.7% binary (.mp4/.png/.wav) — the content scan opened and UTF-8-decoded
# every one of those bytes because nothing skipped binaries. These bounds are
# sized off that measurement: the real text in that workspace was 5.2 MB across
# 1593 files (~3.4 KB average), so none of them fire in the common case.
GREP_BINARY_SNIFF_BYTES = 8192  # same head-sniff size git/ripgrep use
GREP_MAX_FILE_READ_BYTES = 8 * 1024 * 1024  # per-file read cap
GREP_MAX_LINE_CHARS = 4096  # cap fed to regex.search per line/chunk
GREP_SCAN_BYTE_BUDGET = 32 * 1024 * 1024  # total bytes read per call
GREP_SCAN_WALL_CLOCK_SECONDS = 6.0  # total wall clock per call
EXTRACT_ONLY_EXTENSIONS = {
    ".doc", ".docx", ".wps", ".pdf", ".xlsx", ".xls", ".et", ".pptx", ".ppt", ".dps",
}

def _is_user_visible_rel(path: str) -> bool:
    return runtime_user_visible_file_path(path)

# ---------------------------------------------------------------------------
# Schemas
# ---------------------------------------------------------------------------

READ_FILE_SCHEMA = {
    "type": "function",
    "function": {
        "name": "read_file",
        "description": "Read text/images.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "Relative entity path.",
                },
                "offset": {
                    "type": "integer",
                    "description": "0-based line offset.",
                },
                "char_offset": {
                    "type": "integer",
                    "description": "Character offset continuation.",
                },
                "limit": {
                    "type": "integer",
                    "description": "Max lines.",
                },
                "max_chars": {
                    "type": "integer",
                    "description": "Max chars.",
                },
                "expected_sha256": {
                    "type": "string",
                    "description": "Previous source_sha256 guard.",
                },
            },
            "required": ["path"],
        },
    },
}

WRITE_FILE_SCHEMA = {
    "type": "function",
    "function": {
        "name": "write_file",
        "description": "Write/create a user-facing file; use edit_file for small edits.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "Relative entity path.",
                },
                "content": {
                    "type": "string",
                    "description": "Content.",
                },
                "save_to_knowledge": {
                    "type": "boolean",
                    "description": "Override Knowledge sync.",
                },
                "approval_token": {
                    "type": "string",
                    "description": "Approval token.",
                },
                "expected_sha256": {
                    "type": "string",
                    "description": "Previous source_sha256 guard.",
                },
            },
            "required": ["path", "content"],
        },
    },
}

LIST_FILES_SCHEMA = {
    "type": "function",
    "function": {
        "name": "list_files",
        "description": "List raw entity filesystem paths.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "Relative directory.",
                },
                "recursive": {
                    "type": "boolean",
                    "description": "Recursive.",
                },
                "limit": {
                    "type": "integer",
                    "description": "Max entries.",
                },
                "offset": {
                    "type": "integer",
                    "description": "Pagination offset.",
                },
            },
            "required": [],
        },
    },
}

GLOB_FILES_SCHEMA = {
    "type": "function",
    "function": {
        "name": "glob_files",
        "description": "Glob raw entity filesystem paths.",
        "parameters": {
            "type": "object",
            "properties": {
                "pattern": {
                    "type": "string",
                    "description": "Glob pattern.",
                },
                "limit": {
                    "type": "integer",
                    "description": "Max matches.",
                },
                "offset": {
                    "type": "integer",
                    "description": "Pagination offset.",
                },
            },
            "required": ["pattern"],
        },
    },
}

GREP_FILES_SCHEMA = {
    "type": "function",
    "function": {
        "name": "grep_files",
        "description": "Regex search raw entity file contents.",
        "parameters": {
            "type": "object",
            "properties": {
                "pattern": {
                    "type": "string",
                    "description": "Regex pattern.",
                },
                "path": {
                    "type": "string",
                    "description": "Relative directory.",
                },
                "file_glob": {
                    "type": "string",
                    "description": "File glob filter.",
                },
                "max_matches": {
                    "type": "integer",
                    "description": "Max matches.",
                },
                "offset": {
                    "type": "integer",
                    "description": "Match offset.",
                },
                "after": {
                    "type": "string",
                },
            },
            "required": ["pattern"],
        },
    },
}

EDIT_FILE_SCHEMA = {
    "type": "function",
    "function": {
        "name": "edit_file",
        "description": (
            "Edit a file. Text/docx/pptx: old_text/new_text. .xlsx/.xlsm: set_cell, update_row, append_row."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                },
                "operation": {
                    "type": "string",
                    "enum": ["replace_text", "set_cell", "update_row", "append_row"],
                },
                "old_text": {
                    "type": "string",
                },
                "new_text": {
                    "type": "string",
                },
                "replace_all": {
                    "type": "boolean",
                },
                "sheet": {
                    "type": "string",
                },
                "cell": {
                    "type": "string",
                },
                "value": {
                    "type": ["string", "number", "boolean", "null"],
                },
                "header_row": {
                    "type": "integer",
                },
                "match_column": {
                    "type": "string",
                },
                "match_value": {
                    "type": ["string", "number", "boolean", "null"],
                },
                "match_mode": {
                    "type": "string",
                    "enum": ["exact", "contains"],
                },
                "match_case_sensitive": {
                    "type": "boolean",
                },
                "updates": {
                    "type": "object",
                    "additionalProperties": {"type": ["string", "number", "boolean", "null"]},
                },
                "row": {
                    "type": "object",
                    "additionalProperties": {"type": ["string", "number", "boolean", "null"]},
                },
                "values": {
                    "type": "array",
                    "items": {"type": ["string", "number", "boolean", "null"]},
                },
                "approval_token": {
                    "type": "string",
                    "description": "Approval token.",
                },
                "expected_sha256": {
                    "type": "string",
                    "description": "Previous source_sha256 guard.",
                },
            },
            "required": ["path"],
        },
    },
}

DELETE_FILE_SCHEMA = {
    "type": "function",
    "function": {
        "name": "delete_file",
        "description": "Delete a file or empty directory from the entity's filesystem.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "Relative path within entity directory.",
                },
                "approval_token": {
                    "type": "string",
                    "description": "One-time token returned after the user approves deleting a user-visible file.",
                },
                "expected_sha256": {
                    "type": "string",
                    "description": "Optional previous source_sha256; refuses delete if current file changed.",
                },
            },
            "required": ["path"],
        },
    },
}


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _get_entity_root(entity_id: str) -> str | None:
    """Return the entity's filesystem root, or None if FS is disabled."""
    return runtime_entity_file_root(entity_id)


async def _blocked_doc_paths(entity_id: str, rel_paths: list[str], kwargs: dict) -> set[str]:
    """Normalized rel paths that map to a Knowledge document this agent turn's
    user may not read.

    Knowledge documents are real files under the entity root, so the raw file
    tools would otherwise read a private document's bytes without honoring
    ``Document.visibility``. We gate on the same guard the ``/documents`` API
    uses. When the turn has no ``user_id`` (background/system agent) nothing is
    blocked — matching the existing background-caller allowance.
    """
    runtime_context = runtime_tool_call_context_from_kwargs(kwargs)
    user_id = kwargs.get("user_id") or runtime_context.user_id
    if not user_id or not rel_paths:
        return set()
    from packages.core.database import async_session
    from packages.core.services.document_access import unreadable_document_paths

    async with async_session() as db:
        return await unreadable_document_paths(
            db,
            entity_id=entity_id,
            rel_paths=rel_paths,
            user_id=user_id,
            actor_type="agent",
        )


def _safe_path(entity_root: str, relative_path: str) -> str | None:
    """Resolve a path and ensure it stays within entity_root. Returns None if escape detected."""
    root = os.path.realpath(entity_root)
    abs_path = os.path.realpath(os.path.join(root, relative_path))
    if os.path.commonpath([root, abs_path]) != root:
        return None
    return abs_path


async def _workspace_scoped_new_file_path(
    *,
    entity_id: str,
    entity_root: str,
    workspace_id: str | None,
    task_id: str | None = None,
    path: str,
    expected_sha256: str | None = None,
) -> str:
    """Route new workspace file writes under the workspace artifact folder."""
    if not workspace_id:
        return path
    if str(expected_sha256 or "").strip():
        return path

    original_abs = _safe_path(entity_root, path)
    if not original_abs:
        return path

    rel_path = runtime_normalize_entity_file_path(path)
    if not rel_path or not runtime_user_visible_file_path(rel_path):
        return path
    if os.path.exists(original_abs):
        return rel_path

    from packages.core.services.generated_media_naming import (
        resolve_workspace_artifact_base_dir,
        scope_workspace_artifact_path,
    )

    workspace_base_dir = await resolve_workspace_artifact_base_dir(
        entity_id=entity_id,
        workspace_id=workspace_id,
        task_id=task_id,
    )
    if not workspace_base_dir:
        return rel_path
    scoped = scope_workspace_artifact_path(
        rel_path,
        workspace_base_dir,
        default_subdir=WorkspaceArtifactDir.DOCUMENTS.value,
    )
    return runtime_normalize_entity_file_path(scoped) or rel_path


async def _workspace_scoped_existing_path(
    *,
    entity_id: str,
    entity_root: str,
    workspace_id: str | None,
    task_id: str | None = None,
    path: str,
) -> str | None:
    """Where a prior *write* of this logical name would have been rerouted.

    `_write_file` scopes NEW workspace files under the workspace artifact
    folder; read/edit/delete must consult the SAME location so they address
    the file the write created, not a non-existent literal path (the write-A
    read-B fork). Returns the scoped abs_path if that file exists, else None.
    """
    if not workspace_id:
        return None
    rel = runtime_normalize_entity_file_path(path)
    if not rel or not runtime_user_visible_file_path(rel):
        return None
    from packages.core.services.generated_media_naming import (
        resolve_workspace_artifact_base_dir,
        scope_workspace_artifact_path,
    )

    base = await resolve_workspace_artifact_base_dir(
        entity_id=entity_id, workspace_id=workspace_id, task_id=task_id,
    )
    if not base:
        return None
    scoped = runtime_normalize_entity_file_path(
        scope_workspace_artifact_path(rel, base, default_subdir=WorkspaceArtifactDir.DOCUMENTS.value)
    ) or rel
    scoped_abs = _safe_path(entity_root, scoped)
    if scoped_abs and os.path.isfile(scoped_abs):
        return scoped_abs
    return None


def _same_basename_candidates(entity_root: str, path: str, *, limit: int = 5) -> list[str]:
    """Entity-relative paths of files sharing this basename — so a not-found
    error tells the model where the file actually is instead of dead-ending."""
    target = os.path.basename(str(path or "").replace("\\", "/").rstrip("/"))
    if not target:
        return []
    root = os.path.realpath(entity_root)
    found: list[str] = []
    scanned = 0
    for dirpath, dirnames, filenames in os.walk(root):
        dirnames[:] = [d for d in dirnames if not d.startswith(".")]
        for fn in filenames:
            scanned += 1
            if scanned > 5000:
                return found
            if fn == target:
                rel = os.path.relpath(os.path.join(dirpath, fn), root)
                found.append(rel.replace("\\", "/"))
                if len(found) >= limit:
                    return found
    return found


async def _locate_entity_file(
    *,
    entity_id: str,
    entity_root: str,
    path: str,
    workspace_id: str | None,
    task_id: str | None = None,
) -> tuple[str | None, list[str]]:
    """Resolve an EXISTING file for read/edit/delete through ONE code path.

    Order: the literal path, then the workspace-scoped location a prior write
    may have rerouted it to. On miss, returns same-basename candidates so the
    caller can emit a self-correcting error. Returns (abs_path|None, candidates).
    """
    literal = _safe_path(entity_root, path)
    if literal and os.path.isfile(literal):
        return literal, []
    scoped = await _workspace_scoped_existing_path(
        entity_id=entity_id,
        entity_root=entity_root,
        workspace_id=workspace_id,
        task_id=task_id,
        path=path,
    )
    if scoped:
        return scoped, []
    return None, _same_basename_candidates(entity_root, path)


def _not_found_result(path: str, candidates: list[str]) -> str:
    payload: dict[str, Any] = {"error": f"File not found: {path}"}
    if candidates:
        payload["candidates"] = candidates
        payload["hint"] = (
            "No file at that exact path, but a file with the same name exists "
            "elsewhere. Retry with one of `candidates` (these are the real paths)."
        )
    return json.dumps(payload)


_FS_DISABLED_MSG = json.dumps({
    "error": "Entity filesystem is not enabled. Set MANOR_FS_ENABLED=true and MANOR_FS_ROOT.",
})


def _bounded_int(value: Any, default: int, maximum: int, minimum: int = 0) -> int:
    try:
        parsed = int(value)
    except (TypeError, ValueError):
        parsed = default
    return max(minimum, min(parsed, maximum))


def _text_sha256(text: str) -> str:
    return hashlib.sha256(text.encode("utf-8", errors="replace")).hexdigest()


def _file_meta(abs_path: str, content: str) -> dict[str, Any]:
    stat = os.stat(abs_path)
    return {
        "size": stat.st_size,
        "mtime_ns": stat.st_mtime_ns,
        "source_sha256": _text_sha256(content),
    }


def _extract_docx_text(abs_path: str) -> str:
    """Extract plain text from a .docx file."""
    try:
        from docx import Document
        doc = Document(abs_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    paragraphs.append(" | ".join(cells))
        return "\n".join(paragraphs)
    except ImportError:
        # Fallback: extract raw XML text via zipfile
        import zipfile
        import xml.etree.ElementTree as ET
        with zipfile.ZipFile(abs_path, "r") as z:
            with z.open("word/document.xml") as f:
                tree = ET.parse(f)
        ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
        texts = [t.text for t in tree.iter(f"{{{ns['w']}}}t") if t.text]
        return "\n".join(texts)


async def _read_supported_text(abs_path: str) -> str:
    ext = os.path.splitext(abs_path)[1].lower()
    if ext == ".docx":
        return _extract_docx_text(abs_path)
    if ext in {".doc", ".wps", ".pdf", ".xlsx", ".xls", ".et", ".pptx", ".ppt", ".dps"}:
        from packages.core.services.text_extraction import extract_text
        return await extract_text(abs_path, file_type=ext.lstrip("."))
    with open(abs_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def _spreadsheet_json_value(value: Any) -> Any:
    if value is None or isinstance(value, (str, int, float, bool)):
        return value
    return json.dumps(value, ensure_ascii=False, default=str)


def _normalize_sheet_key(value: Any) -> str:
    return re.sub(r"\s+", " ", str(value or "").strip()).lower()


def _spreadsheet_text(value: Any, *, case_sensitive: bool = False) -> str:
    text = str(value if value is not None else "").strip()
    return text if case_sensitive else text.lower()


def _spreadsheet_values_match(
    actual: Any,
    expected: Any,
    *,
    mode: str,
    case_sensitive: bool,
) -> bool:
    actual_text = _spreadsheet_text(actual, case_sensitive=case_sensitive)
    expected_text = _spreadsheet_text(expected, case_sensitive=case_sensitive)
    if mode == "contains":
        return expected_text in actual_text
    return actual_text == expected_text


def _spreadsheet_sheet(workbook, sheet_name: str | None):
    if sheet_name:
        if sheet_name in workbook.sheetnames:
            return workbook[sheet_name]
        lowered = sheet_name.strip().lower()
        for candidate in workbook.sheetnames:
            if candidate.lower() == lowered:
                return workbook[candidate]
        raise ValueError(f"Sheet not found: {sheet_name}")
    return workbook.active


def _spreadsheet_header_map(ws, header_row: int) -> dict[str, int]:
    if header_row < 1:
        raise ValueError("header_row must be >= 1")
    header_map: dict[str, int] = {}
    for cell in ws[header_row]:
        key = _normalize_sheet_key(cell.value)
        if key and key not in header_map:
            header_map[key] = int(cell.column)
    if not header_map:
        raise ValueError(f"No headers found on row {header_row}")
    return header_map


def _spreadsheet_column_index(header_map: dict[str, int], column_name: str) -> int:
    key = _normalize_sheet_key(column_name)
    if not key:
        raise ValueError("Column name is required")
    if key not in header_map:
        known = ", ".join(sorted(header_map.keys())[:20])
        raise ValueError(f"Column not found: {column_name}. Known columns: {known}")
    return header_map[key]


def _replace_text_limited(
    text: str,
    old_text: str,
    new_text: str,
    remaining: int | None,
) -> tuple[str, int, int | None]:
    if old_text not in text:
        return text, 0, remaining
    available = text.count(old_text)
    if remaining is None:
        return text.replace(old_text, new_text), available, None
    replacement_count = min(available, remaining)
    return text.replace(old_text, new_text, replacement_count), replacement_count, remaining - replacement_count


def _set_rich_paragraph_text(paragraph: Any, text: str) -> None:
    try:
        paragraph.text = text
        return
    except Exception:
        pass
    if hasattr(paragraph, "clear"):
        paragraph.clear()
    if hasattr(paragraph, "add_run"):
        run = paragraph.add_run()
        run.text = text


def _replace_in_rich_paragraph(
    paragraph: Any,
    old_text: str,
    new_text: str,
    remaining: int | None,
) -> tuple[int, int | None]:
    replacements = 0
    for run in getattr(paragraph, "runs", []):
        run_text = str(getattr(run, "text", "") or "")
        updated, count, remaining = _replace_text_limited(run_text, old_text, new_text, remaining)
        if count:
            run.text = updated
            replacements += count
            if remaining == 0:
                return replacements, remaining

    paragraph_text = str(getattr(paragraph, "text", "") or "")
    if old_text in paragraph_text:
        updated, count, remaining = _replace_text_limited(paragraph_text, old_text, new_text, remaining)
        if count:
            _set_rich_paragraph_text(paragraph, updated)
            replacements += count
    return replacements, remaining


def _iter_docx_table_paragraphs(tables: Any):
    for table in tables:
        for row in table.rows:
            for cell in row.cells:
                yield from cell.paragraphs
                yield from _iter_docx_table_paragraphs(cell.tables)


def _iter_docx_paragraphs(doc: Any):
    yield from doc.paragraphs
    yield from _iter_docx_table_paragraphs(doc.tables)
    for section in doc.sections:
        for part_name in (
            "header",
            "footer",
            "first_page_header",
            "first_page_footer",
            "even_page_header",
            "even_page_footer",
        ):
            part = getattr(section, part_name, None)
            if part is None:
                continue
            yield from part.paragraphs
            yield from _iter_docx_table_paragraphs(part.tables)


def _replace_docx_sync(abs_path: str, old_text: str, new_text: str, replace_all: bool) -> dict[str, Any]:
    from docx import Document

    doc = Document(abs_path)
    remaining: int | None = None if replace_all else 1
    replacements = 0
    for paragraph in _iter_docx_paragraphs(doc):
        count, remaining = _replace_in_rich_paragraph(paragraph, old_text, new_text, remaining)
        replacements += count
        if remaining == 0:
            break
    if replacements <= 0:
        return {"error": "old_text not found in file. Ensure exact match including whitespace."}
    output = io.BytesIO()
    doc.save(output)
    return {
        "edited": True,
        "operation": "replace_text",
        "file_type": "docx",
        "replacements": replacements,
        "_persisted_bytes": output.getvalue(),
    }


def _iter_pptx_shapes(shapes: Any):
    for shape in shapes:
        yield shape
        child_shapes = getattr(shape, "shapes", None)
        if child_shapes is not None:
            yield from _iter_pptx_shapes(child_shapes)


def _iter_pptx_paragraphs(presentation: Any):
    for slide in presentation.slides:
        for shape in _iter_pptx_shapes(slide.shapes):
            if getattr(shape, "has_text_frame", False):
                yield from shape.text_frame.paragraphs
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        yield from cell.text_frame.paragraphs


def _replace_pptx_sync(abs_path: str, old_text: str, new_text: str, replace_all: bool) -> dict[str, Any]:
    from pptx import Presentation

    presentation = Presentation(abs_path)
    remaining: int | None = None if replace_all else 1
    replacements = 0
    for paragraph in _iter_pptx_paragraphs(presentation):
        count, remaining = _replace_in_rich_paragraph(paragraph, old_text, new_text, remaining)
        replacements += count
        if remaining == 0:
            break
    if replacements <= 0:
        return {"error": "old_text not found in file. Ensure exact match including whitespace."}
    output = io.BytesIO()
    presentation.save(output)
    return {
        "edited": True,
        "operation": "replace_text",
        "file_type": "pptx",
        "replacements": replacements,
        "_persisted_bytes": output.getvalue(),
    }


def _spreadsheet_edit_sync(abs_path: str, params: dict[str, Any]) -> dict[str, Any]:
    from openpyxl import load_workbook
    from openpyxl.utils.cell import coordinate_to_tuple, get_column_letter

    ext = os.path.splitext(abs_path)[1].lower()
    if ext in {".xls", ".et"}:
        return {
            "error": "unsupported_spreadsheet_format",
            "hint": "Convert legacy .xls/.et files to .xlsx before editing.",
        }
    if ext not in {".xlsx", ".xlsm"}:
        return {"error": "not_a_spreadsheet", "hint": "Spreadsheet edit operations only support .xlsx/.xlsm files."}

    keep_vba = ext == ".xlsm"
    wb = load_workbook(abs_path, keep_vba=keep_vba)
    try:
        sheet_name = str(params.get("sheet") or "").strip() or None
        ws = _spreadsheet_sheet(wb, sheet_name)
        operation = str(params.get("operation") or "").strip().lower()
        updated_cells: dict[str, dict[str, Any]] = {}

        if operation == "set_cell":
            cell_ref = str(params.get("cell") or "").strip()
            if not cell_ref:
                return {"error": "cell is required for set_cell"}
            try:
                coordinate_to_tuple(cell_ref)
            except Exception:
                return {"error": f"Invalid cell reference: {cell_ref}"}
            cell = ws[cell_ref]
            old_value = cell.value
            new_value = _spreadsheet_json_value(params.get("value"))
            cell.value = new_value
            updated_cells[cell.coordinate] = {"old": old_value, "new": new_value}
            result: dict[str, Any] = {
                "updated": True,
                "operation": operation,
                "sheet": ws.title,
                "cell": cell.coordinate,
                "updated_cells": updated_cells,
            }

        elif operation == "update_row":
            header_row = _bounded_int(params.get("header_row"), 1, 1000, 1)
            updates = params.get("updates")
            if not isinstance(updates, dict) or not updates:
                return {"error": "updates object is required for update_row"}
            header_map = _spreadsheet_header_map(ws, header_row)
            match_col = _spreadsheet_column_index(header_map, str(params.get("match_column") or ""))
            match_value = params.get("match_value")
            match_mode = str(params.get("match_mode") or "exact").strip().lower()
            if match_mode not in {"exact", "contains"}:
                return {"error": "match_mode must be exact or contains"}
            case_sensitive = bool(params.get("match_case_sensitive"))
            matches: list[int] = []
            for row_idx in range(header_row + 1, ws.max_row + 1):
                if _spreadsheet_values_match(
                    ws.cell(row=row_idx, column=match_col).value,
                    match_value,
                    mode=match_mode,
                    case_sensitive=case_sensitive,
                ):
                    matches.append(row_idx)
            if not matches:
                return {
                    "error": "row_not_found",
                    "match_column": params.get("match_column"),
                    "match_value": match_value,
                }
            if len(matches) > 1:
                return {
                    "error": "multiple_rows_matched",
                    "row_numbers": matches[:50],
                    "hint": "Use a more specific match_value or match_column before updating.",
                }
            row_idx = matches[0]
            for column_name, raw_value in updates.items():
                col_idx = _spreadsheet_column_index(header_map, str(column_name))
                cell = ws.cell(row=row_idx, column=col_idx)
                old_value = cell.value
                new_value = _spreadsheet_json_value(raw_value)
                cell.value = new_value
                updated_cells[cell.coordinate] = {"old": old_value, "new": new_value}
            result = {
                "updated": True,
                "operation": operation,
                "sheet": ws.title,
                "row_number": row_idx,
                "updated_cells": updated_cells,
            }

        elif operation == "append_row":
            row_object = params.get("row")
            values = params.get("values")
            next_row = ws.max_row + 1
            if isinstance(row_object, dict) and row_object:
                header_row = _bounded_int(params.get("header_row"), 1, 1000, 1)
                header_map = _spreadsheet_header_map(ws, header_row)
                for column_name, raw_value in row_object.items():
                    col_idx = _spreadsheet_column_index(header_map, str(column_name))
                    cell = ws.cell(row=next_row, column=col_idx)
                    new_value = _spreadsheet_json_value(raw_value)
                    cell.value = new_value
                    updated_cells[cell.coordinate] = {"old": None, "new": new_value}
            elif isinstance(values, list):
                for col_idx, raw_value in enumerate(values, start=1):
                    cell = ws.cell(row=next_row, column=col_idx)
                    new_value = _spreadsheet_json_value(raw_value)
                    cell.value = new_value
                    updated_cells[cell.coordinate] = {"old": None, "new": new_value}
            else:
                return {"error": "row or values is required for append_row"}
            result = {
                "updated": True,
                "operation": operation,
                "sheet": ws.title,
                "row_number": next_row,
                "updated_cells": updated_cells,
            }

        else:
            return {"error": "operation must be set_cell, update_row, or append_row"}

        output = io.BytesIO()
        wb.save(output)
        result["range"] = ", ".join(
            f"{get_column_letter(coordinate_to_tuple(cell)[1])}{coordinate_to_tuple(cell)[0]}"
            for cell in updated_cells
        )
        result["_persisted_bytes"] = output.getvalue()
        return result
    finally:
        wb.close()


async def _guard_expected_source_sha(
    *,
    abs_path: str,
    path: str,
    expected_sha256: str,
) -> str | None:
    expected_sha256 = (expected_sha256 or "").strip()
    if not expected_sha256:
        return None
    if not os.path.isfile(abs_path):
        return json.dumps({
            "error": "source_missing",
            "path": path,
            "expected_sha256": expected_sha256,
            "hint": "The expected source file no longer exists; re-read or recreate intentionally.",
        })

    content = await _read_supported_text(abs_path)
    meta = _file_meta(abs_path, content)
    if expected_sha256 != meta["source_sha256"]:
        return json.dumps({
            "error": "source_changed",
            "path": path,
            "expected_sha256": expected_sha256,
            "source_sha256": meta["source_sha256"],
            "size": meta["size"],
            "mtime_ns": meta["mtime_ns"],
            "hint": "The file changed since it was read; read it again before editing or overwriting.",
        })
    return None


# ---------------------------------------------------------------------------
# Handlers
# ---------------------------------------------------------------------------

async def _read_file(entity_id: str, **kwargs: Any) -> str:
    runtime_context = runtime_tool_call_context_from_kwargs(kwargs)
    root = _get_entity_root(entity_id)
    if not root:
        return _FS_DISABLED_MSG

    requested_path = str(kwargs.get("path", "") or "")
    if _safe_path(root, requested_path) is None:
        return json.dumps({"error": "Path traversal detected"})
    abs_path, candidates = await _locate_entity_file(
        entity_id=entity_id, entity_root=root, path=requested_path,
        workspace_id=runtime_context.workspace_id,
        task_id=runtime_context.task_id,
    )
    if not abs_path:
        return _not_found_result(requested_path, candidates)

    resolved_path = os.path.relpath(abs_path, root).replace(os.sep, "/")
    if await _blocked_doc_paths(entity_id, [resolved_path], kwargs):
        # Private Knowledge document the caller cannot read: report as missing
        # so the tool never confirms that the document exists.
        return json.dumps({"error": f"File not found: {requested_path}"})

    image_mime_type = {
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".webp": "image/webp",
    }.get(os.path.splitext(abs_path)[1].lower())
    if image_mime_type:
        try:
            with open(abs_path, "rb") as image_file:
                image_bytes = image_file.read(READ_FILE_MAX_IMAGE_BYTES + 1)
            if len(image_bytes) > READ_FILE_MAX_IMAGE_BYTES:
                return json.dumps({
                    "error": "image_too_large",
                    "path": requested_path,
                    "max_bytes": READ_FILE_MAX_IMAGE_BYTES,
                })
            stat = os.stat(abs_path)
            source_sha256 = hashlib.sha256(image_bytes).hexdigest()
            expected_sha256 = str(kwargs.get("expected_sha256") or "").strip()
            if expected_sha256 and expected_sha256 != source_sha256:
                return json.dumps({
                    "error": "source_changed",
                    "path": requested_path,
                    "expected_sha256": expected_sha256,
                    "source_sha256": source_sha256,
                    "size": stat.st_size,
                    "mtime_ns": stat.st_mtime_ns,
                    "hint": "The image changed since it was read; read it again before relying on it.",
                })
            data_url = (
                f"data:{image_mime_type};base64,"
                + base64.b64encode(image_bytes).decode("ascii")
            )
            return json.dumps({
                "path": requested_path,
                "resolved_path": resolved_path,
                "size": stat.st_size,
                "mtime_ns": stat.st_mtime_ns,
                "source_sha256": source_sha256,
                "image": {
                    "data_url": data_url,
                    "mime_type": image_mime_type,
                    "bytes": len(image_bytes),
                    "delivery": "ephemeral_multimodal",
                },
            })
        except Exception as exc:  # noqa: BLE001
            return json.dumps({"error": str(exc)})

    offset = _bounded_int(kwargs.get("offset"), 0, 1_000_000, 0)
    char_offset_arg = kwargs.get("char_offset")
    char_offset = (
        _bounded_int(char_offset_arg, 0, 100_000_000, 0)
        if char_offset_arg is not None
        else None
    )
    limit = _bounded_int(kwargs.get("limit"), READ_FILE_DEFAULT_LINES, READ_FILE_MAX_LINES, 1)
    max_chars = _bounded_int(kwargs.get("max_chars"), READ_FILE_DEFAULT_CHARS, READ_FILE_MAX_CHARS, 1_000)

    try:
        content = await _read_supported_text(abs_path)
        meta = _file_meta(abs_path, content)
        expected_sha256 = str(kwargs.get("expected_sha256") or "").strip()
        if expected_sha256 and expected_sha256 != meta["source_sha256"]:
            return json.dumps({
                "error": "source_changed",
                "path": requested_path,
                "expected_sha256": expected_sha256,
                "source_sha256": meta["source_sha256"],
                "size": meta["size"],
                "mtime_ns": meta["mtime_ns"],
                "hint": "The file changed since the previous slice; restart from offset=0.",
            })

        lines = content.splitlines(keepends=True)
        total_lines = len(lines)
        total_chars = len(content)
        mode = "line"
        start_char_offset = 0
        next_char_offset = None
        partial_line = False
        if char_offset is not None:
            mode = "char"
            start_char_offset = min(char_offset, total_chars)
            content_out = content[start_char_offset : start_char_offset + max_chars]
            selected = content_out.splitlines(keepends=True)
            truncated_by_chars = start_char_offset + len(content_out) < total_chars
            next_offset = None
            next_char_offset = (
                start_char_offset + len(content_out)
                if truncated_by_chars
                else None
            )
        else:
            start_char_offset = sum(len(line) for line in lines[:offset])
            selected = []
            chars_used = 0
            truncated_by_chars = False
            partial = ""
            for line in lines[offset : offset + limit]:
                if chars_used + len(line) <= max_chars:
                    selected.append(line)
                    chars_used += len(line)
                    continue
                truncated_by_chars = True
                if not selected:
                    partial = line[:max_chars]
                    chars_used = len(partial)
                    partial_line = True
                break
            content_out = partial if partial_line else "".join(selected)
            if partial_line:
                selected = []
                next_offset = None
            else:
                next_offset = offset + len(selected) if offset + len(selected) < total_lines else None
            if truncated_by_chars:
                next_char_offset = start_char_offset + len(content_out)
        hint = None
        if truncated_by_chars and next_char_offset is not None:
            hint = (
                "Content hit max_chars; call read_file with char_offset=next_char_offset "
                "and expected_sha256=source_sha256 to continue exactly, or use a "
                "smaller limit/higher max_chars."
            )
        elif truncated_by_chars:
            hint = "Content hit max_chars; use a smaller limit or higher max_chars."
        elif next_offset is not None:
            hint = "Call read_file again with offset=next_offset and expected_sha256=source_sha256 to continue."
        elif next_char_offset is not None:
            hint = "Call read_file again with char_offset=next_char_offset and expected_sha256=source_sha256 to continue."

        return json.dumps({
            "path": requested_path,
            "resolved_path": resolved_path,
            "size": meta["size"],
            "mtime_ns": meta["mtime_ns"],
            "source_sha256": meta["source_sha256"],
            "slice_sha256": _text_sha256(content_out),
            "total_lines": total_lines,
            "total_chars": total_chars,
            "mode": mode,
            "offset": offset,
            "char_offset": start_char_offset,
            "lines_returned": len(selected),
            "next_offset": next_offset,
            "next_char_offset": next_char_offset,
            "partial_line": partial_line,
            "truncated": truncated_by_chars or next_offset is not None,
            "char_truncated": truncated_by_chars,
            "line_truncated": next_offset is not None,
            "max_chars": max_chars,
            "hint": hint,
            "content": content_out,
        })
    except Exception as e:
        return json.dumps({"error": str(e)})


async def _write_file(entity_id: str, **kwargs: Any) -> str:
    runtime_context = runtime_tool_call_context_from_kwargs(kwargs)
    root = _get_entity_root(entity_id)
    if not root:
        return _FS_DISABLED_MSG

    requested_path = str(kwargs.get("path", "") or "")
    path = await _workspace_scoped_new_file_path(
        entity_id=entity_id,
        entity_root=root,
        workspace_id=runtime_context.workspace_id,
        task_id=runtime_context.task_id,
        path=requested_path,
        expected_sha256=str(kwargs.get("expected_sha256") or ""),
    )
    content = kwargs.get("content", "")
    save_flag = kwargs.get("save_to_knowledge", None)
    abs_path = _safe_path(root, path)
    if not abs_path:
        return json.dumps({"error": "Path traversal detected"})

    try:
        stale = await _guard_expected_source_sha(
            abs_path=abs_path,
            path=path,
            expected_sha256=str(kwargs.get("expected_sha256") or ""),
        )
        if stale:
            return stale

        blocked = await runtime_guard_file_mutation(
            entity_id=entity_id,
            user_id=kwargs.get("user_id") or runtime_context.user_id,
            conversation_id=runtime_context.conversation_id,
            tool_name="write_file",
            action="write",
            paths=[path],
            approval_token=kwargs.get("approval_token"),
            content_preview=content,
        )
        if blocked:
            return blocked

        ext = os.path.splitext(path)[1].lower()
        persisted_bytes: bytes
        # Auto-generate real binary for office formats from text content
        if ext == ".pptx":
            from packages.core.services.docgen_service import generate_pptx
            # Extract title from first heading or filename
            lines = content.strip().split("\n")
            title = os.path.splitext(os.path.basename(path))[0]
            for line in lines:
                if line.startswith("# ") and not line.startswith("## "):
                    title = line[2:].strip()
                    break
            persisted_bytes = await generate_pptx(title, content)
        elif ext == ".docx":
            from packages.core.services.docgen_service import generate_docx
            persisted_bytes = await generate_docx(
                os.path.splitext(os.path.basename(path))[0], content
            )
        else:
            persisted_bytes = str(content).encode("utf-8")
        abs_path = runtime_write_entity_file_atomic(
            entity_id,
            path,
            persisted_bytes,
            expected_size=len(persisted_bytes),
            allow_empty=True,
        )

        size = os.path.getsize(abs_path)
        written_content = await _read_supported_text(abs_path)
        written_meta = _file_meta(abs_path, written_content)
        sync = await runtime_sync_entity_file_to_knowledge(
            entity_id=entity_id,
            abs_path=abs_path,
            entity_root=root,
            source="agent",
            created_by=kwargs.get("user_id") or runtime_context.user_id or "ai-agent",
            force=save_flag,
            workspace_id=runtime_context.workspace_id,
            task_id=runtime_context.task_id,
            agent_id=kwargs.get("agent_id") or runtime_context.agent_id,
            conversation_id=runtime_context.conversation_id,
            user_id=kwargs.get("user_id") or runtime_context.user_id,
            tool_name="write_file",
        )

        return json.dumps({
            "written": True,
            "path": path,
            "size": size,
            "source_sha256": written_meta["source_sha256"],
            "mtime_ns": written_meta["mtime_ns"],
            "knowledge_synced": sync.synced,
            "document_id": sync.document_id,
            "knowledge_sync_reason": sync.reason,
        })
    except Exception as e:
        return json.dumps({"error": str(e)})


def _scan_list_files(
    root: str, abs_path: str, recursive: bool, limit: int, offset: int,
) -> tuple[list[dict], int, bool]:
    """Blocking directory listing. Runs in a worker thread — see the note on
    ``_scan_grep_files``."""
    entries: list[dict] = []
    total_seen = 0
    has_more = False
    if recursive:
        for dirpath, dirnames, filenames in os.walk(abs_path):
            dirnames.sort()
            for fn in sorted(filenames):
                full = os.path.join(dirpath, fn)
                rel = os.path.relpath(full, root)
                total_seen += 1
                if total_seen <= offset:
                    continue
                if len(entries) >= limit:
                    has_more = True
                    break
                entries.append({
                    "path": rel,
                    "type": "file",
                    "size": os.path.getsize(full),
                })
            if has_more:
                break
    else:
        items = sorted(os.listdir(abs_path))
        total_seen = len(items)
        for item in items[offset : offset + limit]:
            rel = os.path.relpath(os.path.join(abs_path, item), root)
            full = os.path.join(abs_path, item)
            entries.append({
                "path": rel,
                "type": "dir" if os.path.isdir(full) else "file",
                "size": os.path.getsize(full) if os.path.isfile(full) else None,
            })
        has_more = offset + len(entries) < total_seen
    return entries, total_seen, has_more


async def _list_files(entity_id: str, **kwargs: Any) -> str:
    root = _get_entity_root(entity_id)
    if not root:
        return _FS_DISABLED_MSG

    rel_path = kwargs.get("path", "")
    recursive = bool(kwargs.get("recursive", False))
    limit = _bounded_int(kwargs.get("limit"), LIST_FILES_DEFAULT_LIMIT, LIST_FILES_MAX_LIMIT, 1)
    offset = _bounded_int(kwargs.get("offset"), 0, 100_000, 0)

    abs_path = _safe_path(root, rel_path) if rel_path else root
    if not abs_path:
        return json.dumps({"error": "Path traversal detected"})

    if not os.path.isdir(abs_path):
        return json.dumps({"error": f"Directory not found: {rel_path}"})

    try:
        entries, total_seen, has_more = await asyncio.to_thread(
            _scan_list_files, root, abs_path, recursive, limit, offset,
        )
    except Exception as e:
        return json.dumps({"error": str(e)})

    blocked = await _blocked_doc_paths(
        entity_id, [e["path"] for e in entries if e.get("type") == "file"], kwargs
    )
    if blocked:
        from packages.core.services.knowledge_visibility import normalize_rel_path
        entries = [e for e in entries if normalize_rel_path(e["path"]) not in blocked]

    return json.dumps({
        "count": len(entries),
        "limit": limit,
        "offset": offset,
        "next_offset": offset + len(entries) if has_more else None,
        "has_more": has_more,
        "total": total_seen if not recursive else None,
        "entries": entries,
    })


def _scan_glob_files(
    root: str, pattern: str, limit: int, offset: int,
) -> tuple[list[str], bool]:
    """Blocking name-only tree walk. Runs in a worker thread — see the note on
    ``_scan_grep_files``."""
    matches: list[str] = []
    matched_seen = 0
    has_more = False
    for dirpath, dirnames, filenames in os.walk(root):
        dirnames.sort()
        for fn in sorted(filenames):
            full = os.path.join(dirpath, fn)
            rel = os.path.relpath(full, root)
            if fnmatch.fnmatch(rel, pattern):
                matched_seen += 1
                if matched_seen <= offset:
                    continue
                if len(matches) >= limit:
                    has_more = True
                    break
                matches.append(rel)
        if has_more:
            break
    return matches, has_more


async def _glob_files(entity_id: str, **kwargs: Any) -> str:
    root = _get_entity_root(entity_id)
    if not root:
        return _FS_DISABLED_MSG

    pattern = kwargs.get("pattern", "")
    if not pattern:
        return json.dumps({"error": "pattern is required"})

    limit = _bounded_int(kwargs.get("limit"), GLOB_FILES_DEFAULT_LIMIT, LIST_FILES_MAX_LIMIT, 1)
    offset = _bounded_int(kwargs.get("offset"), 0, 100_000, 0)
    try:
        matches, has_more = await asyncio.to_thread(
            _scan_glob_files, root, pattern, limit, offset,
        )
    except Exception as e:
        return json.dumps({"error": str(e)})

    blocked = await _blocked_doc_paths(entity_id, matches, kwargs)
    if blocked:
        from packages.core.services.knowledge_visibility import normalize_rel_path
        matches = [m for m in matches if normalize_rel_path(m) not in blocked]

    return json.dumps({
        "pattern": pattern,
        "count": len(matches),
        "limit": limit,
        "offset": offset,
        "next_offset": offset + len(matches) if has_more else None,
        "has_more": has_more,
        "files": matches,
    })


def _grep_is_binary(full: str, sniff_bytes: int = GREP_BINARY_SNIFF_BYTES) -> bool:
    """Same head-sniff git/ripgrep use: a NUL byte in the first chunk means
    binary. An unreadable file is treated as skip, not as a scan failure."""
    try:
        with open(full, "rb") as fh:
            return b"\x00" in fh.read(sniff_bytes)
    except OSError:
        return True


def _scan_grep_files(
    root: str,
    search_root: str,
    regex: "re.Pattern[str]",
    file_glob: str,
    max_matches: int,
    offset: int,
    after: str | None = None,
) -> dict:
    """Blocking content scan — MUST NOT run on the event loop (see #424).

    Production incident: this walked a 1.9 GB / 2410-file workspace, opening and
    UTF-8 decoding every file (1.4 GB of it .mp4), and took 8m48s. Measured on
    that same workspace, 99.7% of the bytes were binary and the real text was
    5.2 MB across 1593 files — so the binary skip below removes almost the
    entire cost, and the byte/file/wall-clock bounds are sized generously
    around that real text volume, not around the binary noise.

    ``after`` resumes a budget-truncated scan: it is the exact relative path a
    prior call returned as ``resume_cursor`` (the last file it fully scanned).
    We replay the SAME walk (deterministic: sorted dirnames/filenames) and skip
    every file until we see that exact path again — comparing path strings
    directly, not re-deriving a lexicographic order, since os.walk's traversal
    (all of a directory's own files before any subdirectory) is NOT plain
    lexicographic order over full relative paths. Skipping does no file I/O, so
    replaying up to the cursor is cheap regardless of how far in it is.

    Residual risk, stated plainly rather than papered over: capping bytes fed to
    a single regex.search call (``GREP_MAX_LINE_CHARS``) bounds the quadratic
    blowup this incident's pattern shape could produce (``.*x.*y.*``-style), but
    CPython's ``re`` has no timeout and cannot be interrupted mid-match, so a
    genuinely pathological pattern (nested quantifiers, e.g. ``(a+)+$``) is not
    fully bounded by truncation at any practical cap size. What actually
    contains that case is #424: this now runs in a worker thread, so a hung
    match burns one thread instead of freezing the event loop for every user.
    """
    start = time.monotonic()
    matches: list[dict] = []
    matched_seen = 0
    has_more = False
    files_scanned = 0
    files_skipped_binary = 0
    files_truncated = 0
    bytes_scanned = 0
    truncated = False
    resume_cursor: str | None = None
    # The file the budget check last saw fully handled (scanned OR skipped as
    # binary) — NOT the file the budget trips on. If resume_cursor named the
    # file we were about to look at when the budget ran out, resuming with
    # `after` set to it would skip that exact file via the cursor-replay
    # `continue` below and it would never get scanned at all.
    last_covered_rel: str | None = None
    cursor_found = after is None
    stop = False

    for dirpath, dirnames, filenames in os.walk(search_root):
        dirnames.sort()
        for fn in sorted(filenames):
            if file_glob and not fnmatch.fnmatch(fn, file_glob):
                continue

            full = os.path.join(dirpath, fn)
            rel = os.path.relpath(full, root)

            if not cursor_found:
                if rel == after:
                    cursor_found = True
                continue

            if (
                bytes_scanned >= GREP_SCAN_BYTE_BUDGET
                or (time.monotonic() - start) >= GREP_SCAN_WALL_CLOCK_SECONDS
            ):
                truncated = True
                resume_cursor = last_covered_rel
                stop = True
                break

            if _grep_is_binary(full):
                files_skipped_binary += 1
                last_covered_rel = rel
                continue

            try:
                with open(full, "rb") as fh:
                    raw = fh.read(GREP_MAX_FILE_READ_BYTES + 1)
            except OSError:
                last_covered_rel = rel
                continue

            if len(raw) > GREP_MAX_FILE_READ_BYTES:
                raw = raw[:GREP_MAX_FILE_READ_BYTES]
                files_truncated += 1

            files_scanned += 1
            bytes_scanned += len(raw)
            last_covered_rel = rel
            text = raw.decode("utf-8", errors="replace")

            for line_num, line in enumerate(text.split("\n"), 1):
                line = line[:GREP_MAX_LINE_CHARS]
                if regex.search(line):
                    matched_seen += 1
                    if matched_seen <= offset:
                        continue
                    if len(matches) >= max_matches:
                        has_more = True
                        break
                    matches.append({
                        "file": rel,
                        "line": line_num,
                        "text": line.rstrip()[:500],
                    })
            if has_more:
                break
        if has_more or stop:
            break

    return {
        "matches": matches,
        "has_more": has_more,
        "files_scanned": files_scanned,
        "files_skipped_binary": files_skipped_binary,
        "files_truncated": files_truncated,
        "bytes_scanned": bytes_scanned,
        "truncated": truncated,
        "resume_cursor": resume_cursor,
        "cursor_found": cursor_found,
    }


async def _grep_files(entity_id: str, **kwargs: Any) -> str:
    root = _get_entity_root(entity_id)
    if not root:
        return _FS_DISABLED_MSG

    pattern_str = kwargs.get("pattern", "")
    if not pattern_str:
        return json.dumps({"error": "pattern is required"})

    rel_path = kwargs.get("path", "")
    file_glob = kwargs.get("file_glob", "")
    max_matches = _bounded_int(kwargs.get("max_matches"), GREP_FILES_DEFAULT_LIMIT, LIST_FILES_MAX_LIMIT, 1)
    offset = _bounded_int(kwargs.get("offset"), 0, 100_000, 0)
    after = str(kwargs.get("after") or "").strip() or None

    search_root = _safe_path(root, rel_path) if rel_path else root
    if not search_root or not os.path.isdir(search_root):
        return json.dumps({"error": f"Directory not found: {rel_path}"})

    try:
        regex = re.compile(pattern_str, re.IGNORECASE)
    except re.error as e:
        return json.dumps({"error": f"Invalid regex: {e}"})

    try:
        result = await asyncio.to_thread(
            _scan_grep_files, root, search_root, regex, file_glob, max_matches, offset, after,
        )
        if after and not result["cursor_found"]:
            # The tree changed since the cursor was issued (file moved/deleted)
            # and replay never found it. Skipping the rest of the tree in that
            # state would silently report "no more matches" when we actually
            # never looked — rescan from the start instead of trusting a stale
            # cursor. Rare: a full rescan is cheap once binaries are skipped.
            result = await asyncio.to_thread(
                _scan_grep_files, root, search_root, regex, file_glob, max_matches, offset, None,
            )
    except Exception as e:
        return json.dumps({"error": str(e)})

    matches = result["matches"]
    has_more = result["has_more"]
    blocked = await _blocked_doc_paths(
        entity_id, [m["file"] for m in matches], kwargs
    )
    if blocked:
        from packages.core.services.knowledge_visibility import normalize_rel_path
        matches = [m for m in matches if normalize_rel_path(m["file"]) not in blocked]

    payload: dict[str, Any] = {
        "pattern": pattern_str,
        "count": len(matches),
        "limit": max_matches,
        "offset": offset,
        "next_offset": offset + len(matches) if has_more else None,
        "has_more": has_more,
        "matches": matches,
        "files_scanned": result["files_scanned"],
        "files_skipped_binary": result["files_skipped_binary"],
        "files_truncated": result["files_truncated"],
        "bytes_scanned": result["bytes_scanned"],
        "truncated": result["truncated"],
        "resume_cursor": result["resume_cursor"],
    }
    if result["truncated"]:
        # Keeping this out of GREP_FILES_SCHEMA (a fixed per-conversation token
        # cost, budget-tested) and putting it here instead — it only costs
        # anything on the rare truncated response, not on every conversation.
        payload["hint"] = (
            "Scan ran out of time/budget before covering every file. Retry "
            "with after=resume_cursor to continue; do not report no match yet."
        )
    return json.dumps(payload)


async def _edit_file(entity_id: str, **kwargs: Any) -> str:
    runtime_context = runtime_tool_call_context_from_kwargs(kwargs)
    root = _get_entity_root(entity_id)
    if not root:
        return _FS_DISABLED_MSG

    path = kwargs.get("path", "")
    if _safe_path(root, path) is None:
        return json.dumps({"error": "Path traversal detected"})
    abs_path, candidates = await _locate_entity_file(
        entity_id=entity_id, entity_root=root, path=path,
        workspace_id=runtime_context.workspace_id,
        task_id=runtime_context.task_id,
    )
    if not abs_path:
        return _not_found_result(path, candidates)
    # Address the file we actually located (a workspace-scoped write target),
    # so the sync/knowledge write below updates the SAME projection.
    path = os.path.relpath(abs_path, root).replace("\\", "/")

    try:
        ext = os.path.splitext(abs_path)[1].lower()
        operation = str(kwargs.get("operation") or "").strip().lower()
        spreadsheet_operations = {"set_cell", "update_row", "append_row"}
        if ext in {".xlsx", ".xlsm", ".xls", ".et"}:
            if operation in spreadsheet_operations:
                spreadsheet_kwargs = dict(kwargs)
                spreadsheet_kwargs["_source_tool_name"] = "edit_file"
                return await _edit_spreadsheet(entity_id, **spreadsheet_kwargs)
            return json.dumps({
                "error": "unsupported_binary_edit",
                "path": path,
                "hint": (
                    "Spreadsheet files cannot be edited with text replacement. For .xlsx/.xlsm, "
                    "call edit_file with operation set_cell, update_row, or append_row. Convert "
                    "legacy .xls/.et files to .xlsx before editing."
                ),
            })
        if ext in {".docx", ".pptx"}:
            if operation and operation != "replace_text":
                return json.dumps({
                    "error": "unsupported_operation_for_file_type",
                    "path": path,
                    "operation": operation,
                    "hint": "Use operation replace_text or omit operation for Word/PowerPoint text edits.",
                })
            old_text = kwargs.get("old_text", "")
            new_text = kwargs.get("new_text", "")
            replace_all = bool(kwargs.get("replace_all", False))
            if not old_text:
                return json.dumps({"error": "old_text is required"})

            stale = await _guard_expected_source_sha(
                abs_path=abs_path,
                path=path,
                expected_sha256=str(kwargs.get("expected_sha256") or ""),
            )
            if stale:
                return stale

            blocked = await runtime_guard_file_mutation(
                entity_id=entity_id,
                user_id=kwargs.get("user_id") or kwargs.get("_user_id_from_context"),
                conversation_id=kwargs.get("conversation_id"),
                tool_name="edit_file",
                action="edit",
                paths=[path],
                approval_token=kwargs.get("approval_token"),
                content_preview={
                    "path": path,
                    "operation": "replace_text",
                    "replace_all": replace_all,
                    "old_text": old_text,
                    "new_text": new_text,
                },
            )
            if blocked:
                return blocked

            replace_sync = _replace_docx_sync if ext == ".docx" else _replace_pptx_sync
            result = await asyncio.to_thread(replace_sync, abs_path, old_text, new_text, replace_all)
            if result.get("error"):
                return json.dumps(result, ensure_ascii=False)
            persisted_bytes = result.pop("_persisted_bytes", None)
            if not isinstance(persisted_bytes, bytes):
                return json.dumps({"error": "edited file bytes were not produced"}, ensure_ascii=False)
            abs_path = runtime_write_entity_file_atomic(
                entity_id,
                path,
                persisted_bytes,
                expected_size=len(persisted_bytes),
                allow_empty=False,
            )

            sync = await runtime_sync_entity_file_to_knowledge(
                entity_id=entity_id,
                abs_path=abs_path,
                entity_root=root,
                source="agent",
                created_by=kwargs.get("user_id") or kwargs.get("_user_id_from_context") or "ai-agent",
                force=True,
                workspace_id=kwargs.get("workspace_id"),
                task_id=kwargs.get("task_id"),
                agent_id=kwargs.get("agent_id") or kwargs.get("_agent_id_from_context"),
                conversation_id=kwargs.get("conversation_id"),
                user_id=kwargs.get("user_id") or kwargs.get("_user_id_from_context"),
                tool_name="edit_file",
            )
            content = await _read_supported_text(abs_path)
            meta = _file_meta(abs_path, content)
            result.update({
                "path": path,
                "size": meta["size"],
                "source_sha256": meta["source_sha256"],
                "mtime_ns": meta["mtime_ns"],
                "knowledge_synced": sync.synced,
                "document_id": sync.document_id,
                "knowledge_sync_reason": sync.reason,
            })
            return json.dumps(result, ensure_ascii=False)
        if ext in EXTRACT_ONLY_EXTENSIONS:
            return json.dumps({
                "error": "unsupported_binary_edit",
                "path": path,
                "hint": (
                    "This file type is read through structured extraction and cannot be safely "
                    "edited with text replacement. Export it to Markdown/CSV for text updates, or "
                    "use a file-type-specific edit operation when available."
                ),
            })
        if operation and operation != "replace_text":
            return json.dumps({
                "error": "unsupported_operation_for_file_type",
                "path": path,
                "operation": operation,
                "hint": "Use operation replace_text or omit operation for plain text files.",
            })

        old_text = kwargs.get("old_text", "")
        new_text = kwargs.get("new_text", "")
        replace_all = bool(kwargs.get("replace_all", False))
        if not old_text:
            return json.dumps({"error": "old_text is required"})

        stale = await _guard_expected_source_sha(
            abs_path=abs_path,
            path=path,
            expected_sha256=str(kwargs.get("expected_sha256") or ""),
        )
        if stale:
            return stale

        blocked = await runtime_guard_file_mutation(
            entity_id=entity_id,
            user_id=kwargs.get("user_id") or runtime_context.user_id,
            conversation_id=runtime_context.conversation_id,
            tool_name="edit_file",
            action="edit",
            paths=[path],
            approval_token=kwargs.get("approval_token"),
            content_preview={
                "path": path,
                "operation": operation or "replace_text",
                "replace_all": replace_all,
                "old_text": old_text,
                "new_text": new_text,
            },
        )
        if blocked:
            return blocked

        content = await _read_supported_text(abs_path)

        if old_text not in content:
            return json.dumps({"error": "old_text not found in file. Ensure exact match including whitespace."})

        if replace_all:
            count = content.count(old_text)
            new_content = content.replace(old_text, new_text)
        else:
            count = 1
            new_content = content.replace(old_text, new_text, 1)

        data = new_content.encode("utf-8")
        abs_path = runtime_write_entity_file_atomic(
            entity_id,
            path,
            data,
            expected_size=len(data),
            allow_empty=True,
        )
        sync = await runtime_sync_entity_file_to_knowledge(
            entity_id=entity_id,
            abs_path=abs_path,
            entity_root=root,
            source="agent",
            created_by=kwargs.get("user_id") or runtime_context.user_id or "ai-agent",
            force=True,
            workspace_id=runtime_context.workspace_id,
            task_id=runtime_context.task_id,
            agent_id=kwargs.get("agent_id") or runtime_context.agent_id,
            conversation_id=runtime_context.conversation_id,
            user_id=kwargs.get("user_id") or runtime_context.user_id,
            tool_name="edit_file",
        )
        meta = _file_meta(abs_path, new_content)

        return json.dumps({
            "edited": True,
            "path": path,
            "replacements": count,
            "source_sha256": _text_sha256(new_content),
            "size": meta["size"],
            "mtime_ns": meta["mtime_ns"],
            "knowledge_synced": sync.synced,
            "document_id": sync.document_id,
            "knowledge_sync_reason": sync.reason,
        })
    except Exception as e:
        return json.dumps({"error": str(e)})


async def _edit_spreadsheet(entity_id: str, **kwargs: Any) -> str:
    root = _get_entity_root(entity_id)
    if not root:
        return _FS_DISABLED_MSG

    path = kwargs.get("path", "")
    tool_name = str(kwargs.get("_source_tool_name") or "edit_spreadsheet")
    abs_path = _safe_path(root, path)
    if not abs_path:
        return json.dumps({"error": "Path traversal detected"})
    if not os.path.isfile(abs_path):
        return json.dumps({"error": f"File not found: {path}"})

    try:
        stale = await _guard_expected_source_sha(
            abs_path=abs_path,
            path=path,
            expected_sha256=str(kwargs.get("expected_sha256") or ""),
        )
        if stale:
            return stale

        preview = {
            "path": path,
            "sheet": kwargs.get("sheet"),
            "operation": kwargs.get("operation"),
            "cell": kwargs.get("cell"),
            "match_column": kwargs.get("match_column"),
            "match_value": kwargs.get("match_value"),
            "updates": kwargs.get("updates"),
            "row": kwargs.get("row"),
            "values": kwargs.get("values"),
        }
        blocked = await runtime_guard_file_mutation(
            entity_id=entity_id,
            user_id=kwargs.get("user_id") or kwargs.get("_user_id_from_context"),
            conversation_id=kwargs.get("conversation_id"),
            tool_name=tool_name,
            action="edit",
            paths=[path],
            approval_token=kwargs.get("approval_token"),
            content_preview=preview,
        )
        if blocked:
            return blocked

        result = await asyncio.to_thread(_spreadsheet_edit_sync, abs_path, dict(kwargs))
        if result.get("error"):
            return json.dumps(result, ensure_ascii=False)
        persisted_bytes = result.pop("_persisted_bytes", None)
        if not isinstance(persisted_bytes, bytes):
            return json.dumps({"error": "edited spreadsheet bytes were not produced"}, ensure_ascii=False)
        abs_path = runtime_write_entity_file_atomic(
            entity_id,
            path,
            persisted_bytes,
            expected_size=len(persisted_bytes),
            allow_empty=False,
        )

        sync = await runtime_sync_entity_file_to_knowledge(
            entity_id=entity_id,
            abs_path=abs_path,
            entity_root=root,
            source="agent",
            created_by=kwargs.get("user_id") or kwargs.get("_user_id_from_context") or "ai-agent",
            force=True,
            workspace_id=kwargs.get("workspace_id"),
            task_id=kwargs.get("task_id"),
            agent_id=kwargs.get("agent_id") or kwargs.get("_agent_id_from_context"),
            conversation_id=kwargs.get("conversation_id"),
            user_id=kwargs.get("user_id") or kwargs.get("_user_id_from_context"),
            tool_name=tool_name,
        )
        content = await _read_supported_text(abs_path)
        meta = _file_meta(abs_path, content)
        result.update({
            "path": path,
            "size": meta["size"],
            "source_sha256": meta["source_sha256"],
            "mtime_ns": meta["mtime_ns"],
            "knowledge_synced": sync.synced,
            "document_id": sync.document_id,
            "knowledge_sync_reason": sync.reason,
        })
        return json.dumps(result, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False)


async def _delete_file(entity_id: str, **kwargs: Any) -> str:
    runtime_context = runtime_tool_call_context_from_kwargs(kwargs)
    root = _get_entity_root(entity_id)
    if not root:
        return _FS_DISABLED_MSG

    path = kwargs.get("path", "")
    abs_path = _safe_path(root, path)
    if not abs_path:
        return json.dumps({"error": "Path traversal detected"})

    if not os.path.exists(abs_path):
        # A literal miss may still resolve to a workspace-scoped file (dirs
        # only ever match literally).
        located, candidates = await _locate_entity_file(
            entity_id=entity_id, entity_root=root, path=path,
            workspace_id=runtime_context.workspace_id,
            task_id=runtime_context.task_id,
        )
        if not located:
            return _not_found_result(path, candidates)
        abs_path = located

    try:
        rel_for_trash = os.path.relpath(abs_path, root)
        expected_sha256 = str(kwargs.get("expected_sha256") or "")
        if expected_sha256:
            if not os.path.isfile(abs_path):
                return json.dumps({
                    "error": "source_not_file",
                    "path": rel_for_trash,
                    "hint": "expected_sha256 can only guard file deletes, not directories.",
                })
            stale = await _guard_expected_source_sha(
                abs_path=abs_path,
                path=rel_for_trash,
                expected_sha256=expected_sha256,
            )
            if stale:
                return stale

        blocked = await runtime_guard_file_mutation(
            entity_id=entity_id,
            user_id=kwargs.get("user_id") or runtime_context.user_id,
            conversation_id=runtime_context.conversation_id,
            tool_name="delete_file",
            action="delete",
            paths=[rel_for_trash],
            approval_token=kwargs.get("approval_token"),
            content_preview={"delete": rel_for_trash},
        )
        if blocked:
            return blocked

        if os.path.isfile(abs_path):
            os.remove(abs_path)
            if _is_user_visible_rel(rel_for_trash):
                await runtime_trash_knowledge_path(entity_id, rel_for_trash)
            return json.dumps({"deleted": True, "path": path, "type": "file"})
        elif os.path.isdir(abs_path):
            if os.listdir(abs_path):
                return json.dumps({"error": "Directory not empty. Remove contents first."})
            os.rmdir(abs_path)
            if _is_user_visible_rel(rel_for_trash):
                await runtime_trash_knowledge_path(entity_id, rel_for_trash)
            return json.dumps({"deleted": True, "path": path, "type": "directory"})
        else:
            return json.dumps({"error": f"Unknown path type: {path}"})
    except Exception as e:
        return json.dumps({"error": str(e)})


# ---------------------------------------------------------------------------
# Export
# ---------------------------------------------------------------------------

def get_tools() -> list[tuple[dict, callable]]:
    return [
        (READ_FILE_SCHEMA, _read_file),
        (WRITE_FILE_SCHEMA, _write_file),
        (EDIT_FILE_SCHEMA, _edit_file),
        (DELETE_FILE_SCHEMA, _delete_file),
        (LIST_FILES_SCHEMA, _list_files),
        (GLOB_FILES_SCHEMA, _glob_files),
        (GREP_FILES_SCHEMA, _grep_files),
    ]
